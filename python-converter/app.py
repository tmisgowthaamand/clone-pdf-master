#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flask API for PowerPoint to PDF Conversion
Professional-grade conversion using LibreOffice
"""

import sys
import io
import os

# Fix Windows encoding issues - force UTF-8
if sys.platform == 'win32':
    # Set console output encoding to UTF-8
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
    # Set environment variable for subprocess encoding
    os.environ['PYTHONIOENCODING'] = 'utf-8'

from flask import Flask, request, send_file, jsonify, session
from flask_cors import CORS
import os
import tempfile
from pathlib import Path
import subprocess
import shutil
from werkzeug.utils import secure_filename
import pandas as pd
import camelot
import base64
import uuid
import json
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.lib.colors import HexColor
from datetime import datetime
import requests

app = Flask(__name__)
app.secret_key = os.urandom(24)  # For session management

# Configure CORS to allow frontend access
ALLOWED_ORIGINS = os.environ.get('ALLOWED_ORIGINS', '*').split(',')
# Add Vercel frontend domain explicitly
if ALLOWED_ORIGINS == ['*']:
    ALLOWED_ORIGINS = [
        'http://localhost:5173',
        'http://localhost:3000',
        'https://pdf-tools-phi.vercel.app',
        'https://*.vercel.app'
    ]

CORS(app, resources={
    r"/api/*": {
        "origins": ALLOWED_ORIGINS,
        "methods": ["GET", "POST", "OPTIONS", "PUT", "DELETE"],
        "allow_headers": ["Content-Type", "Authorization", "Accept"],
        "supports_credentials": True,
        "expose_headers": ["Content-Disposition"]
    }
})

# Add CORS headers to all responses
@app.after_request
def after_request(response):
    origin = request.headers.get('Origin')
    if origin:
        # Check if origin is allowed
        if '*' in ALLOWED_ORIGINS or origin in ALLOWED_ORIGINS or any(origin.endswith(o.replace('https://*.', '.')) for o in ALLOWED_ORIGINS if '*' in o):
            response.headers['Access-Control-Allow-Origin'] = origin
            response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS, PUT, DELETE'
            response.headers['Access-Control-Allow-Headers'] = 'Content-Type, Authorization, Accept'
            response.headers['Access-Control-Expose-Headers'] = 'Content-Disposition'
            response.headers['Access-Control-Allow-Credentials'] = 'true'
    return response

# Add LibreOffice to PATH (cross-platform)
if sys.platform == 'win32':
    libreoffice_path = r"C:\Program Files\LibreOffice\program"
    if libreoffice_path not in os.environ.get('PATH', ''):
        os.environ['PATH'] = libreoffice_path + os.pathsep + os.environ.get('PATH', '')

ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'pdf', 'doc', 'docx'}
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB

def allowed_file(filename, allowed_types=None):
    if allowed_types is None:
        allowed_types = ALLOWED_EXTENSIONS
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_types

def convert_with_libreoffice(input_path, output_dir):
    """Convert using LibreOffice - Optimized for speed and reliability"""
    # Detect LibreOffice executable based on platform
    if sys.platform == 'win32':
        soffice_exe = r"C:\Program Files\LibreOffice\program\soffice.exe"
        # Kill any existing LibreOffice processes to avoid file locks (Windows only)
        try:
            subprocess.run(['taskkill', '/F', '/IM', 'soffice.exe', '/T'], 
                          capture_output=True, timeout=5, encoding='utf-8', errors='replace')
            subprocess.run(['taskkill', '/F', '/IM', 'soffice.bin', '/T'], 
                          capture_output=True, timeout=5, encoding='utf-8', errors='replace')
            import time
            time.sleep(0.5)  # Reduced wait time
        except:
            pass
    else:
        # Linux/Unix - use system LibreOffice
        soffice_exe = shutil.which('soffice') or 'soffice'
    
    # Optimized command with faster settings
    cmd = [
        soffice_exe,
        '--headless',
        '--invisible',
        '--nologo',
        '--nofirststartwizard',
        '--norestore',
        '--nolockcheck',  # Skip lock check for faster startup
        '--convert-to', 'pdf',
        '--outdir', str(output_dir),
        str(input_path)
    ]
    
    print(f"Converting with LibreOffice: {Path(input_path).name}")
    
    # Reduced timeout to 60 seconds for faster failure detection
    result = subprocess.run(
        cmd, 
        capture_output=True, 
        text=True, 
        timeout=60, 
        encoding='utf-8', 
        errors='replace'
    )
    
    if result.returncode != 0:
        print(f"LibreOffice error: {result.stderr}")
        raise RuntimeError(f"Conversion failed: {result.stderr}")
    
    pdf_name = Path(input_path).stem + '.pdf'
    pdf_path = Path(output_dir) / pdf_name
    
    if not pdf_path.exists():
        raise RuntimeError("PDF was not created")
    
    print(f"Conversion complete: {pdf_name}")
    return pdf_path

def convert_pdf_to_pptx(input_path, output_dir):
    """
    Convert PDF to PPTX by converting pages to images - Optimized for speed
    Uses pdf2image and python-pptx
    """
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile as tmp
        
        print(f"Converting PDF to PowerPoint: {input_path}")
        
        # Optimized: Use lower DPI for faster conversion (200 instead of 300)
        # Still maintains good quality while being 2x faster
        dpi = 200
        
        # Convert PDF pages to images (poppler is in PATH on Linux/Docker)
        images = convert_from_path(
            str(input_path),
            dpi=dpi,
            fmt='png',
            thread_count=2  # Use 2 threads for faster processing
        )
        
        print(f"Converted {len(images)} pages to images")
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add each image as a slide with FULL BACKGROUND
        for idx, image in enumerate(images, 1):
            # Save image temporarily
            with tmp.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                # Optimize: Use lower quality PNG for faster save
                image.save(tmp_img.name, 'PNG', optimize=True)
                tmp_img_path = tmp_img.name
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # FULL BACKGROUND: Fill entire slide with image
            slide.shapes.add_picture(
                tmp_img_path,
                0,  # left = 0 (no margin)
                0,  # top = 0 (no margin)
                width=prs.slide_width,   # Full width
                height=prs.slide_height  # Full height
            )
            
            # Clean up immediately
            try:
                os.unlink(tmp_img_path)
            except:
                pass
            
            print(f"Added slide {idx}/{len(images)}")
        
        # Save PPTX
        pptx_name = Path(input_path).stem + '.pptx'
        pptx_path = Path(output_dir) / pptx_name
        prs.save(str(pptx_path))
        
        print(f"PowerPoint created: {pptx_path}")
        return pptx_path
        
    except ImportError as e:
        raise RuntimeError(
            f"Missing required libraries: {str(e)}"
        )
    except Exception as e:
        error_msg = str(e)
        if "poppler" in error_msg.lower() or "pdftoppm" in error_msg.lower():
            raise RuntimeError(
                "Poppler not found. Install poppler-utils package."
            )
        raise RuntimeError(f"PDF to PPTX conversion failed: {error_msg}")

def convert_pdf_to_docx(input_path, output_dir):
    """
    Convert PDF to DOCX using pdf2docx - Best free Python solution
    Preserves text, tables, basic images, fonts, and layout
    """
    cv = None
    try:
        print(f"\n{'='*60}")
        print(f"PDF TO DOCX CONVERSION - pdf2docx")
        print(f"Input: {input_path}")
        print(f"Output dir: {output_dir}")
        print(f"{'='*60}\n")
        
        from pdf2docx import Converter
        
        # Output path
        docx_name = Path(input_path).stem + '.docx'
        docx_path = Path(output_dir) / docx_name
        
        print(f"Output path: {docx_path}")
        
        # Create converter instance
        cv = Converter(str(input_path))
        print("[OK] Converter instance created")
        
        # Convert PDF to DOCX
        print("Starting conversion...")
        cv.convert(str(docx_path), start=0, end=None)
        print("[OK] Conversion completed")
        
        # Close converter
        cv.close()
        print("[OK] Converter closed")
        
        if not docx_path.exists():
            raise RuntimeError("DOCX file was not created")
        
        print(f"[OK] DOCX file created successfully: {docx_path}")
        print(f"{'='*60}\n")
        return docx_path
        
    except ImportError as e:
        print(f"Import error: {str(e)}")
        if cv:
            try:
                cv.close()
            except:
                pass
        raise RuntimeError(
            f"Missing required library. Install with: pip install pdf2docx. Error: {str(e)}"
        )
    except Exception as e:
        print(f"Conversion error: {str(e)}")
        import traceback
        traceback.print_exc()
        if cv:
            try:
                cv.close()
            except:
                pass
        raise RuntimeError(f"PDF to DOCX conversion failed: {str(e)}")

@app.route('/api/convert/pptx-to-pdf', methods=['POST'])
def convert_pptx_to_pdf():
    """Convert PowerPoint to PDF - Optimized for speed"""
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Only .ppt and .pptx allowed'}), 400
    
    tmpdir = None
    try:
        print(f"Converting PowerPoint to PDF: {file.filename}")
        
        # Create temporary directory with unique name
        import uuid
        tmpdir = tempfile.mkdtemp(prefix=f'pptx_convert_{uuid.uuid4().hex[:8]}_')
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        input_path = Path(tmpdir) / filename
        file.save(str(input_path))
        
        # Convert to PDF using optimized LibreOffice
        pdf_path = convert_with_libreoffice(input_path, tmpdir)
        
        # Read PDF into memory before cleanup
        with open(pdf_path, 'rb') as f:
            pdf_data = f.read()
        
        print(f"PowerPoint to PDF conversion complete: {len(pdf_data)} bytes")
        
        # Clean up temp directory immediately
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except:
            pass
        
        # Send PDF from memory
        from io import BytesIO
        return send_file(
            BytesIO(pdf_data),
            mimetype='application/pdf',
            as_attachment=True,
            download_name=filename.rsplit('.', 1)[0] + '.pdf'
        )
            
    except Exception as e:
        print(f"PowerPoint to PDF error: {str(e)}")
        # Clean up on error
        if tmpdir and os.path.exists(tmpdir):
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass
        return jsonify({'error': str(e)}), 500

@app.route('/api/convert/pdf-to-pptx', methods=['POST'])
def convert_pdf_to_pptx_endpoint():
    """Convert PDF to PowerPoint endpoint"""
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename, {'pdf'}):
        return jsonify({'error': 'Invalid file type. Only .pdf allowed'}), 400
    
    tmpdir = None
    try:
        # Create temporary directory with unique name
        import uuid
        tmpdir = tempfile.mkdtemp(prefix=f'pdf_convert_{uuid.uuid4().hex[:8]}_')
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        input_path = Path(tmpdir) / filename
        file.save(str(input_path))
        
        # Convert to PPTX
        pptx_path = convert_pdf_to_pptx(input_path, tmpdir)
        
        # Read PPTX into memory before cleanup
        with open(pptx_path, 'rb') as f:
            pptx_data = f.read()
        
        # Clean up temp directory
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except:
            pass
        
        # Send PPTX from memory
        from io import BytesIO
        return send_file(
            BytesIO(pptx_data),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename.rsplit('.', 1)[0] + '.pptx'
        )
            
    except Exception as e:
        # Clean up on error
        if tmpdir and os.path.exists(tmpdir):
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass
        return jsonify({'error': str(e)}), 500

@app.route('/api/convert/docx-to-pdf', methods=['POST'])
def convert_docx_to_pdf_endpoint():
    """Convert Word to PDF endpoint"""
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename, {'doc', 'docx'}):
        return jsonify({'error': 'Invalid file type. Only .doc and .docx allowed'}), 400
    
    tmpdir = None
    try:
        # Create temporary directory with unique name
        import uuid
        tmpdir = tempfile.mkdtemp(prefix=f'word_pdf_{uuid.uuid4().hex[:8]}_')
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        input_path = Path(tmpdir) / filename
        file.save(str(input_path))
        
        # Convert to PDF using LibreOffice
        pdf_path = convert_with_libreoffice(input_path, tmpdir)
        
        # Read PDF into memory before cleanup
        with open(pdf_path, 'rb') as f:
            pdf_data = f.read()
        
        # Clean up temp directory
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except:
            pass
        
        # Send PDF from memory
        from io import BytesIO
        return send_file(
            BytesIO(pdf_data),
            mimetype='application/pdf',
            as_attachment=True,
            download_name=filename.rsplit('.', 1)[0] + '.pdf'
        )
            
    except Exception as e:
        # Clean up on error
        if tmpdir and os.path.exists(tmpdir):
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass
        return jsonify({'error': str(e)}), 500

@app.route('/api/convert/pdf-to-docx', methods=['POST'])
def convert_pdf_to_docx_endpoint():
    """Convert PDF to Word endpoint"""
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename, {'pdf'}):
        return jsonify({'error': 'Invalid file type. Only .pdf allowed'}), 400
    
    tmpdir = None
    try:
        # Create temporary directory with unique name
        import uuid
        tmpdir = tempfile.mkdtemp(prefix=f'pdf_word_{uuid.uuid4().hex[:8]}_')
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        input_path = Path(tmpdir) / filename
        file.save(str(input_path))
        
        # Convert to DOCX
        docx_path = convert_pdf_to_docx(input_path, tmpdir)
        
        # Read DOCX into memory before cleanup
        with open(docx_path, 'rb') as f:
            docx_data = f.read()
        
        # Clean up temp directory
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except:
            pass
        
        # Send DOCX from memory
        from io import BytesIO
        return send_file(
            BytesIO(docx_data),
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=filename.rsplit('.', 1)[0] + '.docx'
        )
            
    except Exception as e:
        # Clean up on error
        if tmpdir and os.path.exists(tmpdir):
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass
        return jsonify({'error': str(e)}), 500

@app.route('/api/convert/pdf-to-excel', methods=['POST'])
def pdf_to_excel():
    """Convert PDF to Excel - Optimized for speed"""
    tmpdir = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'pdf'}):
            return jsonify({'error': 'Only PDF files are allowed'}), 400
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        # Save uploaded PDF
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(tmpdir, filename)
        file.save(pdf_path)
        
        print(f"Converting PDF to Excel: {filename}")
        
        # Use PyMuPDF to extract images, text, and tables
        import fitz
        doc = fitz.open(pdf_path)
        
        # Create Excel file
        excel_name = Path(filename).stem + '.xlsx'
        excel_path = os.path.join(tmpdir, excel_name)
        
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.drawing.image import Image as OpenpyxlImage
        from openpyxl.utils import get_column_letter
        
        wb = Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        
        current_row = 1
        
        # Process first page to extract logo and headers
        page = doc[0]
        
        # Extract and add images (Bank logo)
        image_list = page.get_images()
        if image_list:
            print(f"Found {len(image_list)} images (logos)")
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image temporarily
                    img_path = os.path.join(tmpdir, f'logo.{image_ext}')
                    with open(img_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    # Add image to Excel (top right, like bank statements)
                    excel_img = OpenpyxlImage(img_path)
                    # Resize to reasonable size
                    if excel_img.height > 80:
                        ratio = 80 / excel_img.height
                        excel_img.height = 80
                        excel_img.width = int(excel_img.width * ratio)
                    
                    # Position in top right (column H)
                    ws.add_image(excel_img, 'H1')
                    print(f"  Added logo to Excel")
                except Exception as e:
                    print(f"  Warning: Could not add image: {e}")
        
        # Extract text for title and headers
        text_dict = page.get_text("dict")
        blocks = text_dict["blocks"]
        
        # Find "Detailed Statement" title
        for block in blocks:
            if block["type"] == 0:  # Text block
                for line in block["lines"]:
                    line_text = ""
                    for span in line["spans"]:
                        line_text += span["text"]
                    if "Detailed Statement" in line_text:
                        # Add title (centered, bold, large)
                        ws.merge_cells('A2:G2')
                        title_cell = ws['A2']
                        title_cell.value = "Detailed Statement"
                        title_cell.font = Font(size=16, bold=True)
                        title_cell.alignment = Alignment(horizontal='center', vertical='center')
                        ws.row_dimensions[2].height = 25  # Make title row taller
                        print("  Added 'Detailed Statement' title")
                        break
        
        # Start table from row 4 (right after title)
        current_row = 4
        
        # Now extract tables using Camelot
        print("\nExtracting tables...")
        all_tables = []
        
        # Try lattice mode first
        try:
            print("Strategy 1: Lattice mode (bordered tables)...")
            tables = camelot.read_pdf(
                str(pdf_path), 
                pages='all', 
                flavor='lattice',
                line_scale=40,
                shift_text=['l', 't']
            )
            if tables and len(tables) > 0:
                for idx, table in enumerate(tables):
                    df = table.df
                    if not df.empty and len(df) > 1:
                        df = df.replace('', pd.NA).dropna(how='all', axis=0).dropna(how='all', axis=1)
                        
                        if len(df) > 0:
                            headers = df.iloc[0].astype(str).str.strip()
                            
                            valid_cols = []
                            for i, col in enumerate(headers):
                                col_data = df.iloc[1:, i]
                                if (col != '' and col != 'nan' and not pd.isna(col)) or col_data.notna().any():
                                    valid_cols.append(i)
                            
                            df = df.iloc[:, valid_cols]
                            headers = df.iloc[0].astype(str).str.strip()
                            
                            seen = {}
                            unique_headers = []
                            for i, col in enumerate(headers):
                                if col == '' or col == 'nan' or pd.isna(col):
                                    col = f'Column_{i+1}'
                                if col in seen:
                                    seen[col] += 1
                                    unique_headers.append(f'{col}_{seen[col]}')
                                else:
                                    seen[col] = 0
                                    unique_headers.append(col)
                            
                            df.columns = unique_headers
                            df = df[1:].reset_index(drop=True)
                            df = df.dropna(how='all')
                            
                            all_tables.append(df)
                            print(f"  [OK] Table {idx + 1}: {len(df)} rows x {len(df.columns)} columns")
        except Exception as e:
            print(f"  Lattice mode failed: {str(e)}")
        
        # Try stream mode if no tables found
        if not all_tables:
            try:
                print("Strategy 2: Stream mode (borderless tables)...")
                tables = camelot.read_pdf(
                    str(pdf_path), 
                    pages='all', 
                    flavor='stream',
                    edge_tol=50,
                    row_tol=10,
                    column_tol=10
                )
                if tables and len(tables) > 0:
                    for idx, table in enumerate(tables):
                        df = table.df
                        if not df.empty and len(df) > 1:
                            df = df.replace('', pd.NA).dropna(how='all', axis=0).dropna(how='all', axis=1)
                            
                            if len(df) > 0:
                                headers = df.iloc[0].astype(str).str.strip()
                                
                                valid_cols = []
                                for i, col in enumerate(headers):
                                    col_data = df.iloc[1:, i]
                                    if (col != '' and col != 'nan' and not pd.isna(col)) or col_data.notna().any():
                                        valid_cols.append(i)
                                
                                df = df.iloc[:, valid_cols]
                                headers = df.iloc[0].astype(str).str.strip()
                                
                                seen = {}
                                unique_headers = []
                                for i, col in enumerate(headers):
                                    if col == '' or col == 'nan' or pd.isna(col):
                                        col = f'Column_{i+1}'
                                    if col in seen:
                                        seen[col] += 1
                                        unique_headers.append(f'{col}_{seen[col]}')
                                    else:
                                        seen[col] = 0
                                        unique_headers.append(col)
                                
                                df.columns = unique_headers
                                df = df[1:].reset_index(drop=True)
                                df = df.dropna(how='all')
                                
                                all_tables.append(df)
                                print(f"  [OK] Table {idx + 1}: {len(df)} rows x {len(df.columns)} columns")
            except Exception as e:
                print(f"  Stream mode failed: {str(e)}")
        
        if not all_tables:
            # If no tables found, still save the workbook with logo and title
            print("Warning: No tables found, saving logo and title only")
            wb.save(excel_path)
            doc.close()
            
            return send_file(
                excel_path,
                as_attachment=True,
                download_name=excel_name,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
        
        # Combine all tables with same structure
        tables_by_cols = {}
        for df in all_tables:
            col_count = len(df.columns)
            if col_count not in tables_by_cols:
                tables_by_cols[col_count] = []
            tables_by_cols[col_count].append(df)
        
        main_col_count = max(tables_by_cols.keys(), key=lambda k: len(tables_by_cols[k]))
        main_tables = tables_by_cols[main_col_count]
        
        print(f"[OK] Found {len(main_tables)} tables with {main_col_count} columns")
        
        if len(main_tables) == 1:
            combined_df = main_tables[0]
        else:
            base_columns = list(main_tables[0].columns)
            normalized_tables = []
            
            for df in main_tables:
                df_copy = df.copy()
                df_copy.columns = base_columns
                normalized_tables.append(df_copy)
            
            combined_df = pd.concat(normalized_tables, ignore_index=True)
        
        print(f"[OK] Combined table: {len(combined_df)} rows x {len(combined_df.columns)} columns")
        
        # Convert NA values to empty strings for Excel compatibility
        combined_df = combined_df.fillna('')
        
        # Write table data to worksheet starting from current_row
        # Write headers
        for col_idx, col_name in enumerate(combined_df.columns, start=1):
            cell = ws.cell(row=current_row, column=col_idx, value=str(col_name))
            cell.font = Font(bold=True, size=11)
            cell.fill = PatternFill(start_color='D3D3D3', end_color='D3D3D3', fill_type='solid')
            cell.alignment = Alignment(horizontal='center', vertical='center')
        
        current_row += 1
        
        # Write data rows
        for _, row_data in combined_df.iterrows():
            for col_idx, value in enumerate(row_data, start=1):
                # Convert value to string if it's not empty
                cell_value = str(value) if value != '' else ''
                cell = ws.cell(row=current_row, column=col_idx, value=cell_value)
                cell.alignment = Alignment(horizontal='center', vertical='center')
            current_row += 1
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = get_column_letter(column[0].column)
            for cell in column:
                try:
                    if cell.value and len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max(max_length + 2, 12), 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # Save workbook
        wb.save(excel_path)
        doc.close()
        
        print(f"[OK] Excel file created with logo and formatting: {excel_path}")
        print(f"{'='*60}\n")
        
        return send_file(
            excel_path,
            as_attachment=True,
            download_name=excel_name,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    
    except Exception as e:
        print(f"\n{'='*60}")
        print(f"ERROR IN PDF TO EXCEL CONVERSION")
        print(f"Error: {str(e)}")
        print(f"{'='*60}")
        import traceback
        traceback.print_exc()
        print(f"{'='*60}\n")
        return jsonify({'error': str(e)}), 500
    finally:
        # Cleanup will happen after file is sent
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

# OLD CODE BELOW - REMOVE
def pdf_to_excel_OLD_BACKUP():
    """OLD VERSION - DO NOT USE"""
    tmpdir = None
    try:
        # Try lattice mode first (for tables with borders) - FASTER
        try:
            print("Strategy 1: Lattice mode (bordered tables)...")
            tables = camelot.read_pdf(
                str(pdf_path), 
                pages='all', 
                flavor='lattice',
                line_scale=40,  # Better line detection
                shift_text=['l', 't']  # Align text to left-top for better accuracy
            )
            if tables and len(tables) > 0:
                for idx, table in enumerate(tables):
                    df = table.df
                    if not df.empty and len(df) > 1:
                        # Clean up the dataframe
                        df = df.replace('', pd.NA).dropna(how='all', axis=0).dropna(how='all', axis=1)
                        
                        if len(df) > 0:
                            # Use first row as header and clean column names
                            headers = df.iloc[0].astype(str).str.strip()
                            
                            # Debug: Print detected headers
                            print(f"    Detected headers: {list(headers)}")
                            
                            # Remove columns where header is empty/nan AND all data is empty
                            valid_cols = []
                            for i, col in enumerate(headers):
                                col_data = df.iloc[1:, i]  # Get column data (excluding header row)
                                # Keep column if header is not empty OR if column has data
                                if (col != '' and col != 'nan' and not pd.isna(col)) or col_data.notna().any():
                                    valid_cols.append(i)
                            
                            # Filter to valid columns only
                            df = df.iloc[:, valid_cols]
                            headers = df.iloc[0].astype(str).str.strip()
                            
                            # Handle duplicate or empty column names
                            seen = {}
                            unique_headers = []
                            for i, col in enumerate(headers):
                                if col == '' or col == 'nan' or pd.isna(col):
                                    col = f'Column_{i+1}'
                                if col in seen:
                                    seen[col] += 1
                                    unique_headers.append(f'{col}_{seen[col]}')
                                else:
                                    seen[col] = 0
                                    unique_headers.append(col)
                            
                            df.columns = unique_headers
                            df = df[1:].reset_index(drop=True)
                            
                            # Remove any completely empty rows
                            df = df.dropna(how='all')
                            
                            all_tables.append(df)
                            print(f"  [OK] Table {idx + 1}: {len(df)} rows x {len(df.columns)} columns (accuracy: {table.accuracy:.1f}%)")
        except Exception as e:
            print(f"  Lattice mode failed: {str(e)}")
        
        # Try stream mode if no tables found
        if not all_tables:
            try:
                print("Strategy 2: Stream mode (borderless tables)...")
                tables = camelot.read_pdf(
                    str(pdf_path), 
                    pages='all', 
                    flavor='stream',
                    edge_tol=50,
                    row_tol=10,
                    column_tol=10
                )
                if tables and len(tables) > 0:
                    for idx, table in enumerate(tables):
                        df = table.df
                        if not df.empty and len(df) > 1:
                            # Clean up the dataframe
                            df = df.replace('', pd.NA).dropna(how='all', axis=0).dropna(how='all', axis=1)
                            
                            if len(df) > 0:
                                # Use first row as header and clean column names
                                headers = df.iloc[0].astype(str).str.strip()
                                
                                # Remove columns where header is empty/nan AND all data is empty
                                valid_cols = []
                                for i, col in enumerate(headers):
                                    col_data = df.iloc[1:, i]  # Get column data (excluding header row)
                                    # Keep column if header is not empty OR if column has data
                                    if (col != '' and col != 'nan' and not pd.isna(col)) or col_data.notna().any():
                                        valid_cols.append(i)
                                
                                # Filter to valid columns only
                                df = df.iloc[:, valid_cols]
                                headers = df.iloc[0].astype(str).str.strip()
                                
                                # Handle duplicate or empty column names
                                seen = {}
                                unique_headers = []
                                for i, col in enumerate(headers):
                                    if col == '' or col == 'nan' or pd.isna(col):
                                        col = f'Column_{i+1}'
                                    if col in seen:
                                        seen[col] += 1
                                        unique_headers.append(f'{col}_{seen[col]}')
                                    else:
                                        seen[col] = 0
                                        unique_headers.append(col)
                                
                                df.columns = unique_headers
                                df = df[1:].reset_index(drop=True)
                                
                                # Remove any completely empty rows
                                df = df.dropna(how='all')
                                
                                all_tables.append(df)
                                print(f"  [OK] Table {idx + 1}: {len(df)} rows x {len(df.columns)} columns")
            except Exception as e:
                print(f"  Stream mode failed: {str(e)}")
        
        if not all_tables:
            return jsonify({'error': 'No tables found in PDF'}), 400
        
        # Combine all tables with the SAME structure (same columns)
        # Group tables by number of columns
        tables_by_cols = {}
        for df in all_tables:
            col_count = len(df.columns)
            if col_count not in tables_by_cols:
                tables_by_cols[col_count] = []
            tables_by_cols[col_count].append(df)
        
        # Use the group with most tables (main data structure)
        main_col_count = max(tables_by_cols.keys(), key=lambda k: len(tables_by_cols[k]))
        main_tables = tables_by_cols[main_col_count]
        
        print(f"[OK] Found {len(main_tables)} tables with {main_col_count} columns")
        print(f"  (Ignored {len(all_tables) - len(main_tables)} tables with different structure)")
        
        # Combine all main tables into one
        if len(main_tables) == 1:
            combined_df = main_tables[0]
        else:
            # Ensure all tables have same column names
            base_columns = list(main_tables[0].columns)
            
            # Debug: Print base columns being used
            print(f"  Using column names: {base_columns}")
            
            normalized_tables = []
            
            for df in main_tables:
                df_copy = df.copy()
                df_copy.columns = base_columns  # Use first table's column names
                normalized_tables.append(df_copy)
            
            combined_df = pd.concat(normalized_tables, ignore_index=True)
        
        print(f"[OK] Combined table: {len(combined_df)} rows x {len(combined_df.columns)} columns")
        
        # Create Excel file
        excel_name = Path(filename).stem + '.xlsx'
        excel_path = os.path.join(tmpdir, excel_name)
        
        from openpyxl.styles import Font, Alignment, PatternFill
        
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            # Write combined table
            combined_df.to_excel(writer, sheet_name='Sheet1', index=False)
            
            # Format worksheet
            worksheet = writer.sheets['Sheet1']
            
            # Style header row (bold, centered, light gray background)
            header_fill = PatternFill(start_color='D3D3D3', end_color='D3D3D3', fill_type='solid')
            header_font = Font(bold=True, size=11)
            center_alignment = Alignment(horizontal='center', vertical='center')
            
            for cell in worksheet[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center_alignment
            
            # Auto-adjust column widths and center align data
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                        # Center align all cells
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    except:
                        pass
                adjusted_width = min(max(max_length + 2, 12), 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
        
        print(f"[OK] Excel file created: {excel_path}")
        print(f"[OK] Extracted {len(all_tables)} tables")
        print(f"{'='*60}\n")
        
        # Send file
        return send_file(
            excel_path,
            as_attachment=True,
            download_name=excel_name,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    
    except Exception as e:
        print(f"\n{'='*60}")
        print(f"ERROR IN PDF TO EXCEL CONVERSION")
        print(f"Error: {str(e)}")
        print(f"{'='*60}")
        import traceback
        traceback.print_exc()
        print(f"{'='*60}\n")
        return jsonify({'error': str(e)}), 500
    finally:
        # Cleanup will happen after file is sent
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

def convert_excel_to_pdf_advanced(excel_path, output_dir):
    """
    Advanced Excel to PDF conversion with perfect formatting preservation
    Matches iLovePDF quality - preserves images, logos, borders, fonts, colors
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.drawing.image import Image as OpenpyxlImage
        from reportlab.lib.pagesizes import A4, letter, landscape
        from reportlab.lib.units import inch, mm
        from reportlab.pdfgen import canvas
        from reportlab.lib import colors
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from PIL import Image as PILImage
        import io
        
        print(f"\n{'='*60}")
        print(f"ADVANCED EXCEL TO PDF - iLovePDF Quality")
        print(f"Input: {excel_path}")
        print(f"{'='*60}\n")
        
        # Load workbook
        wb = load_workbook(excel_path, data_only=True)
        ws = wb.active
        
        # Get dimensions
        max_row = ws.max_row
        max_col = ws.max_column
        
        print(f"Sheet: {ws.title}")
        print(f"Dimensions: {max_row} rows × {max_col} columns")
        
        # Calculate column widths and row heights
        col_widths = {}
        for col_idx in range(1, max_col + 1):
            # Get column letter - handle merged cells
            from openpyxl.utils import get_column_letter
            col_letter = get_column_letter(col_idx)
            col_dim = ws.column_dimensions[col_letter]
            
            # Calculate actual content width for better sizing
            max_content_width = 0
            for row_idx in range(1, min(max_row + 1, 100)):  # Check first 100 rows
                cell = ws.cell(row_idx, col_idx)
                # Skip merged cells for width calculation
                if hasattr(cell, 'value') and cell.value:
                    content_len = len(str(cell.value))
                    font_size = 10
                    if hasattr(cell, 'font') and cell.font and cell.font.size:
                        font_size = cell.font.size
                    # Estimate width: char_count * font_size * 0.6
                    estimated_width = content_len * font_size * 0.6
                    max_content_width = max(max_content_width, estimated_width)
            
            # Use the larger of: Excel width or content width
            excel_width = (col_dim.width or 10) * 7
            width_pts = max(excel_width, max_content_width + 10)  # Add padding
            col_widths[col_idx] = width_pts
        
        row_heights = {}
        for row_idx in range(1, max_row + 1):
            row_dim = ws.row_dimensions[row_idx]
            # Convert Excel height to points
            height_pts = (row_dim.height or 15) * 0.75
            row_heights[row_idx] = height_pts
        
        # Calculate total dimensions
        total_width = sum(col_widths.values())
        total_height = sum(row_heights.values())
        
        # Set page size with margins
        margin = 20
        page_width = total_width + (2 * margin)
        page_height = total_height + (2 * margin)
        
        # Use A4 if dimensions fit, otherwise custom
        if page_width <= 595 and page_height <= 842:  # A4 dimensions
            pagesize = A4
            if page_width > page_height:
                pagesize = landscape(A4)
        else:
            pagesize = (page_width, page_height)
        
        print(f"Page size: {pagesize[0]:.1f} × {pagesize[1]:.1f} pts")
        
        # Create PDF
        pdf_name = Path(excel_path).stem + '.pdf'
        pdf_path = Path(output_dir) / pdf_name
        c = canvas.Canvas(str(pdf_path), pagesize=pagesize)
        
        # Starting position
        x_start = margin
        y_start = pagesize[1] - margin
        
        # Extract and draw images/logos
        images_drawn = 0
        if hasattr(ws, '_images'):
            for img in ws._images:
                try:
                    # Get image data
                    img_data = img._data()
                    pil_img = PILImage.open(io.BytesIO(img_data))
                    
                    # Get anchor position
                    anchor = img.anchor
                    if hasattr(anchor, '_from'):
                        col_idx = anchor._from.col + 1
                        row_idx = anchor._from.row + 1
                        
                        # Calculate position
                        x_pos = x_start + sum(col_widths[i] for i in range(1, col_idx))
                        y_pos = y_start - sum(row_heights[i] for i in range(1, row_idx))
                        
                        # Calculate image size
                        img_width = img.width * 0.75 if hasattr(img, 'width') else pil_img.width * 0.75
                        img_height = img.height * 0.75 if hasattr(img, 'height') else pil_img.height * 0.75
                        
                        # Draw image
                        img_reader = ImageReader(io.BytesIO(img_data))
                        c.drawImage(img_reader, x_pos, y_pos - img_height, 
                                   width=img_width, height=img_height, 
                                   preserveAspectRatio=True, mask='auto')
                        images_drawn += 1
                        print(f"  [OK] Image {images_drawn} drawn at row {row_idx}, col {col_idx}")
                except Exception as e:
                    print(f"  Warning: Could not draw image: {str(e)}")
        
        # Draw cells
        y_pos = y_start
        for row_idx in range(1, max_row + 1):
            x_pos = x_start
            row_height = row_heights[row_idx]
            
            for col_idx in range(1, max_col + 1):
                cell = ws.cell(row_idx, col_idx)
                col_width = col_widths[col_idx]
                
                # Skip merged cells (they'll be handled by the main cell)
                from openpyxl.cell.cell import MergedCell
                if isinstance(cell, MergedCell):
                    x_pos += col_width
                    continue
                
                # Draw cell border
                if hasattr(cell, 'border') and cell.border and (cell.border.left.style or cell.border.right.style or 
                                   cell.border.top.style or cell.border.bottom.style):
                    c.setStrokeColor(colors.black)
                    c.setLineWidth(0.5)
                    
                    if cell.border.left.style:
                        c.line(x_pos, y_pos, x_pos, y_pos - row_height)
                    if cell.border.right.style:
                        c.line(x_pos + col_width, y_pos, x_pos + col_width, y_pos - row_height)
                    if cell.border.top.style:
                        c.line(x_pos, y_pos, x_pos + col_width, y_pos)
                    if cell.border.bottom.style:
                        c.line(x_pos, y_pos - row_height, x_pos + col_width, y_pos - row_height)
                
                # Draw cell background
                if hasattr(cell, 'fill') and cell.fill and hasattr(cell.fill, 'start_color') and cell.fill.start_color and hasattr(cell.fill.start_color, 'rgb') and cell.fill.start_color.rgb:
                    try:
                        rgb = cell.fill.start_color.rgb
                        if len(rgb) == 8:  # ARGB format
                            rgb = rgb[2:]  # Remove alpha
                        r, g, b = int(rgb[0:2], 16)/255, int(rgb[2:4], 16)/255, int(rgb[4:6], 16)/255
                        c.setFillColor(colors.Color(r, g, b))
                        c.rect(x_pos, y_pos - row_height, col_width, row_height, fill=1, stroke=0)
                    except:
                        pass
                
                # Draw cell text
                if cell.value:
                    text = str(cell.value)
                    
                    # Set font
                    font_size = 9  # Slightly smaller for better fit
                    if cell.font and cell.font.size:
                        font_size = min(cell.font.size, 11)  # Cap at 11pt
                    
                    # Set text color
                    text_color = colors.black
                    if cell.font and cell.font.color and cell.font.color.rgb:
                        try:
                            rgb = cell.font.color.rgb
                            if len(rgb) == 8:
                                rgb = rgb[2:]
                            r, g, b = int(rgb[0:2], 16)/255, int(rgb[2:4], 16)/255, int(rgb[4:6], 16)/255
                            text_color = colors.Color(r, g, b)
                        except:
                            pass
                    
                    c.setFillColor(text_color)
                    font_name = "Helvetica-Bold" if (cell.font and cell.font.bold) else "Helvetica"
                    c.setFont(font_name, font_size)
                    
                    # Handle long text - wrap if needed
                    text_width = c.stringWidth(text, font_name, font_size)
                    available_width = col_width - 4  # Leave padding
                    
                    if text_width > available_width and len(text) > 15:
                        # Wrap text for long content
                        words = text.split()
                        lines = []
                        current_line = []
                        
                        for word in words:
                            test_line = ' '.join(current_line + [word])
                            if c.stringWidth(test_line, font_name, font_size) <= available_width:
                                current_line.append(word)
                            else:
                                if current_line:
                                    lines.append(' '.join(current_line))
                                current_line = [word]
                        if current_line:
                            lines.append(' '.join(current_line))
                        
                        # Draw multiple lines
                        line_height = font_size + 2
                        start_y = y_pos - ((row_height - (len(lines) * line_height)) / 2)
                        
                        for i, line in enumerate(lines[:3]):  # Max 3 lines
                            text_y = start_y - (i * line_height)
                            text_x = x_pos + 2
                            
                            if cell.alignment:
                                if cell.alignment.horizontal == 'center':
                                    text_x = x_pos + (col_width / 2)
                                    c.drawCentredString(text_x, text_y, line)
                                elif cell.alignment.horizontal == 'right':
                                    text_x = x_pos + col_width - 2
                                    c.drawRightString(text_x, text_y, line)
                                else:
                                    c.drawString(text_x, text_y, line)
                            else:
                                c.drawString(text_x, text_y, line)
                    else:
                        # Single line text
                        text_x = x_pos + 2
                        text_y = y_pos - (row_height / 2) - (font_size / 3)  # Vertically center
                        
                        if cell.alignment:
                            if cell.alignment.horizontal == 'center':
                                text_x = x_pos + (col_width / 2)
                                c.drawCentredString(text_x, text_y, text)
                            elif cell.alignment.horizontal == 'right':
                                text_x = x_pos + col_width - 2
                                c.drawRightString(text_x, text_y, text)
                            else:
                                c.drawString(text_x, text_y, text)
                        else:
                            c.drawString(text_x, text_y, text)
                
                x_pos += col_width
            
            y_pos -= row_height
        
        # Save PDF
        c.save()
        
        print(f"[OK] PDF created: {pdf_path}")
        print(f"[OK] Images preserved: {images_drawn}")
        print(f"[OK] iLovePDF-quality conversion complete!")
        print(f"{'='*60}\n")
        
        return pdf_path
        
    except ImportError as e:
        print(f"Missing library: {str(e)}")
        print("Falling back to LibreOffice conversion...")
        return None
    except Exception as e:
        print(f"Advanced conversion failed: {str(e)}")
        import traceback
        traceback.print_exc()
        print("Falling back to LibreOffice conversion...")
        return None

def convert_csv_to_excel_api(csv_path, output_dir):
    """Convert CSV to Excel for API"""
    try:
        import pandas as pd
        from openpyxl.styles import Border, Side, Font
        
        # Read CSV with multiple encoding attempts
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        df = None
        
        for encoding in encodings:
            try:
                df = pd.read_csv(csv_path, encoding=encoding)
                break
            except:
                continue
        
        if df is None:
            raise RuntimeError("Could not read CSV file")
        
        # Create Excel file
        base_name = os.path.splitext(os.path.basename(csv_path))[0]
        excel_path = os.path.join(output_dir, f"{base_name}.xlsx")
        
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Sheet1', index=False)
            
            worksheet = writer.sheets['Sheet1']
            
            # Auto-adjust column widths
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
            
            # Add borders
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            for row in worksheet.iter_rows(min_row=1, max_row=len(df)+1, 
                                          min_col=1, max_col=len(df.columns)):
                for cell in row:
                    cell.border = thin_border
            
            # Bold headers
            for cell in worksheet[1]:
                cell.font = Font(bold=True)
        
        return excel_path
    except Exception as e:
        print(f"CSV conversion error: {str(e)}")
        return None

def convert_excel_to_pdf_professional_api(excel_path, output_dir):
    """Professional Excel to PDF using COM automation"""
    try:
        import win32com.client
        import pythoncom
        
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        
        workbook = excel.Workbooks.Open(os.path.abspath(excel_path))
        
        base_name = os.path.splitext(os.path.basename(excel_path))[0]
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        
        # Configure page setup
        for worksheet in workbook.Worksheets:
            worksheet.PageSetup.Zoom = False
            worksheet.PageSetup.FitToPagesWide = 1
            worksheet.PageSetup.FitToPagesTall = False
            
            worksheet.PageSetup.LeftMargin = excel.InchesToPoints(0.2)
            worksheet.PageSetup.RightMargin = excel.InchesToPoints(0.2)
            worksheet.PageSetup.TopMargin = excel.InchesToPoints(0.3)
            worksheet.PageSetup.BottomMargin = excel.InchesToPoints(0.3)
            
            # Landscape for wide tables
            used_range = worksheet.UsedRange
            if used_range.Columns.Count > 6:
                worksheet.PageSetup.Orientation = 2  # Landscape
            
            worksheet.PageSetup.PrintQuality = 600
            worksheet.PageSetup.PaperSize = 9  # A4
            worksheet.PageSetup.CenterHorizontally = True
        
        # Export to PDF
        workbook.ExportAsFixedFormat(
            Type=0,
            Filename=os.path.abspath(pdf_path),
            Quality=0,
            IncludeDocProperties=True,
            IgnorePrintAreas=False,
            OpenAfterPublish=False
        )
        
        workbook.Close(SaveChanges=False)
        excel.Quit()
        pythoncom.CoUninitialize()
        
        return pdf_path
    except:
        return None

@app.route('/api/convert/excel-to-pdf', methods=['POST'])
def excel_to_pdf():
    """Convert Excel/CSV to PDF - iLovePDF quality"""
    tmpdir = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'xlsx', 'xls', 'csv'}):
            return jsonify({'error': 'Only Excel (.xlsx, .xls) or CSV (.csv) files are allowed'}), 400
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        file_path = os.path.join(tmpdir, filename)
        file.save(file_path)
        
        print(f"\n{'='*60}")
        print(f"EXCEL/CSV TO PDF CONVERSION - iLovePDF Quality")
        print(f"Input: {filename}")
        print(f"{'='*60}\n")
        
        # Handle CSV files - convert to Excel first
        excel_path = file_path
        if filename.lower().endswith('.csv'):
            print("Step 1: Converting CSV to Excel...")
            excel_path = convert_csv_to_excel_api(file_path, tmpdir)
            if not excel_path or not os.path.exists(excel_path):
                raise RuntimeError("CSV to Excel conversion failed")
            print(f"[OK] CSV converted to Excel: {os.path.basename(excel_path)}\n")
        
        # Try professional Excel COM conversion first (iLovePDF quality)
        pdf_path = None
        print("Step 2: Converting to PDF with Microsoft Excel COM...")
        pdf_path = convert_excel_to_pdf_professional_api(excel_path, tmpdir)
        
        # Fallback to LibreOffice if COM fails
        if not pdf_path or not os.path.exists(pdf_path):
            print(f"\n{'='*60}")
            print(f"EXCEL TO PDF CONVERSION - LibreOffice (Fallback)")
            print(f"Input: {excel_path}")
            print(f"{'='*60}\n")
            
            # Detect LibreOffice executable based on platform
            if sys.platform == 'win32':
                soffice_exe = r"C:\Program Files\LibreOffice\program\soffice.exe"
                # Kill existing LibreOffice processes (Windows only)
                try:
                    subprocess.run(['taskkill', '/F', '/IM', 'soffice.exe', '/T'], 
                                  capture_output=True, timeout=5, encoding='utf-8', errors='replace')
                    subprocess.run(['taskkill', '/F', '/IM', 'soffice.bin', '/T'], 
                                  capture_output=True, timeout=5, encoding='utf-8', errors='replace')
                    import time
                    time.sleep(1)
                except:
                    pass
            else:
                # Linux/Unix - use system LibreOffice
                soffice_exe = shutil.which('soffice') or 'soffice'
            
            # Convert with optimal settings - preserve images and formatting
            # calc_pdf_Export preserves all images, logos, and formatting
            cmd = [
                soffice_exe,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nofirststartwizard',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf:calc_pdf_Export',
                '--outdir', tmpdir,
                excel_path
            ]
            
            print(f"Running: {' '.join(cmd)}")
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120,
                encoding='utf-8',
                errors='replace'
            )
            
            if result.returncode != 0:
                print(f"LibreOffice error: {result.stderr}")
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
            
            # Find the output PDF
            base_name = os.path.splitext(os.path.basename(excel_path))[0]
            pdf_path = os.path.join(tmpdir, f"{base_name}.pdf")
            
            if not os.path.exists(pdf_path):
                possible_paths = [
                    os.path.join(tmpdir, f"{base_name}.pdf"),
                    os.path.join(tmpdir, f"{filename.rsplit('.', 1)[0]}.pdf")
                ]
                for path in possible_paths:
                    if os.path.exists(path):
                        pdf_path = path
                        break
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise RuntimeError("PDF file was not created")
            
            print(f"[OK] PDF file created: {pdf_path}")
            print(f"{'='*60}\n")
        
        # Send file
        pdf_name = Path(filename).stem + '.pdf'
        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=pdf_name,
            mimetype='application/pdf'
        )
    
    except Exception as e:
        print(f"Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500
    finally:
        # Cleanup
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

@app.route('/api/convert/excel-to-bank-statement', methods=['POST', 'OPTIONS'])
def excel_to_bank_statement():
    """Convert Excel to Bank Statement PDF - Uses professional Excel COM for perfect alignment"""
    # Handle CORS preflight
    if request.method == 'OPTIONS':
        return '', 204
    
    tmpdir = None
    try:
        print(f"\n{'='*60}")
        print(f"EXCEL TO BANK STATEMENT - Request received")
        print(f"Request files: {list(request.files.keys())}")
        print(f"{'='*60}\n")
        
        if 'file' not in request.files:
            print("ERROR: No file in request")
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        print(f"File received: {file.filename}")
        
        if file.filename == '':
            print("ERROR: Empty filename")
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'xlsx', 'xls', 'csv'}):
            print(f"ERROR: Invalid file type: {file.filename}")
            return jsonify({'error': 'Only Excel (.xlsx, .xls) or CSV (.csv) files are allowed'}), 400
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        file_path = os.path.join(tmpdir, filename)
        file.save(file_path)
        
        # Check if logo was uploaded (for future use)
        logo_path = None
        if 'logo' in request.files:
            logo_file = request.files['logo']
            if logo_file and logo_file.filename:
                logo_filename = secure_filename(logo_file.filename)
                logo_path = os.path.join(tmpdir, logo_filename)
                logo_file.save(logo_path)
                print(f"[OK] Logo uploaded: {logo_filename}")
        
        print(f"\n{'='*60}")
        print(f"EXCEL TO BANK STATEMENT PDF - Professional Conversion")
        print(f"Input: {filename}")
        print(f"Logo: {logo_path if logo_path else 'None'}")
        print(f"{'='*60}\n")
        
        # Handle CSV files - convert to Excel first
        excel_path = file_path
        if filename.lower().endswith('.csv'):
            print("Step 1: Converting CSV to Excel...")
            excel_path = convert_csv_to_excel_api(file_path, tmpdir)
            if not excel_path or not os.path.exists(excel_path):
                raise RuntimeError("CSV to Excel conversion failed")
            print(f"[OK] CSV converted to Excel: {os.path.basename(excel_path)}\n")
        
        # Use professional Excel COM conversion for perfect alignment
        pdf_path = None
        print("Step 2: Converting to PDF with Microsoft Excel COM (Perfect Alignment)...")
        pdf_path = convert_excel_to_pdf_professional_api(excel_path, tmpdir)
        
        # Fallback to LibreOffice if COM fails
        if not pdf_path or not os.path.exists(pdf_path):
            print(f"\n{'='*60}")
            print(f"FALLBACK: Using LibreOffice for conversion")
            print(f"{'='*60}\n")
            
            # Detect LibreOffice executable
            if sys.platform == 'win32':
                soffice_exe = r"C:\Program Files\LibreOffice\program\soffice.exe"
                try:
                    subprocess.run(['taskkill', '/F', '/IM', 'soffice.exe', '/T'], 
                                  capture_output=True, timeout=5, encoding='utf-8', errors='replace')
                    subprocess.run(['taskkill', '/F', '/IM', 'soffice.bin', '/T'], 
                                  capture_output=True, timeout=5, encoding='utf-8', errors='replace')
                    import time
                    time.sleep(1)
                except:
                    pass
            else:
                soffice_exe = shutil.which('soffice') or 'soffice'
            
            cmd = [
                soffice_exe,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nofirststartwizard',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf:calc_pdf_Export',
                '--outdir', tmpdir,
                excel_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120,
                                   encoding='utf-8', errors='replace')
            
            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
            
            base_name = os.path.splitext(os.path.basename(excel_path))[0]
            pdf_path = os.path.join(tmpdir, f"{base_name}.pdf")
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise RuntimeError("PDF file was not created")
            
            print(f"[OK] PDF created with LibreOffice")
        
        # Send file
        pdf_name = Path(filename).stem + '_statement.pdf'
        print(f"[OK] Sending PDF: {pdf_name}")
        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=pdf_name,
            mimetype='application/pdf'
        )
    
    except Exception as e:
        error_msg = str(e)
        print(f"\n{'='*60}")
        print(f"ERROR IN EXCEL TO BANK STATEMENT CONVERSION")
        print(f"Error: {error_msg}")
        print(f"{'='*60}")
        import traceback
        traceback.print_exc()
        print(f"{'='*60}\n")
        return jsonify({'error': error_msg}), 500
    finally:
        # Cleanup
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

@app.route('/api/convert/pdf-to-jpg', methods=['POST'])
def pdf_to_jpg():
    """Convert PDF to JPG images - Optimized for speed"""
    tmpdir = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'pdf'}):
            return jsonify({'error': 'Only PDF files are allowed'}), 400
        
        # Get options from form data
        conversion_mode = request.form.get('mode', 'pages')  # 'pages' or 'extract'
        quality = request.form.get('quality', 'high')  # 'low', 'normal', 'high'
        
        # Optimized quality settings - reduced DPI for faster conversion
        quality_map = {
            'low': {'dpi': 100, 'jpeg_quality': 70},
            'normal': {'dpi': 150, 'jpeg_quality': 80},
            'high': {'dpi': 200, 'jpeg_quality': 90}  # Reduced from 300 to 200 for 2x speed
        }
        settings = quality_map.get(quality, quality_map['high'])
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(tmpdir, filename)
        file.save(pdf_path)
        
        print(f"Converting PDF to JPG: {filename} ({quality} quality, {settings['dpi']} DPI)")
        
        import fitz  # PyMuPDF
        from PIL import Image
        import zipfile
        
        doc = fitz.open(pdf_path)
        jpg_paths = []
        
        if conversion_mode == 'extract':
            # Extract embedded images from PDF
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Save extracted image
                    jpg_filename = f"{Path(filename).stem}_page{page_num+1}_img{img_index+1}.jpg"
                    jpg_path = os.path.join(tmpdir, jpg_filename)
                    
                    with open(jpg_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    jpg_paths.append(jpg_path)
                    print(f"[OK] Extracted image {len(jpg_paths)}: {jpg_filename}")
        else:
            # Convert each page to JPG
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=settings['dpi'])
                
                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # Save as JPG
                jpg_filename = f"{Path(filename).stem}_page_{page_num+1}.jpg"
                jpg_path = os.path.join(tmpdir, jpg_filename)
                img.save(jpg_path, "JPEG", quality=settings['jpeg_quality'])
                jpg_paths.append(jpg_path)
                
                print(f"[OK] Converted page {page_num+1}/{len(doc)}: {jpg_filename}")
        
        doc.close()
        
        if not jpg_paths:
            return jsonify({'error': 'No images found or generated'}), 400
        
        # Return list of image URLs for separate downloads
        print(f"[OK] Generated {len(jpg_paths)} images")
        print(f"{'='*60}\n")
        
        # Store files temporarily and return file info
        file_info = []
        for jpg_path in jpg_paths:
            file_info.append({
                'filename': os.path.basename(jpg_path),
                'path': jpg_path
            })
        
        # Store in session or return base64 encoded images
        images_data = []
        for jpg_path in jpg_paths:
            with open(jpg_path, 'rb') as f:
                img_base64 = base64.b64encode(f.read()).decode('utf-8')
                images_data.append({
                    'filename': os.path.basename(jpg_path),
                    'data': img_base64
                })
        
        return jsonify({
            'success': True,
            'count': len(images_data),
            'images': images_data
        })
    
    except Exception as e:
        print(f"Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500
    finally:
        # Cleanup
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

@app.route('/api/convert/jpg-to-pdf', methods=['POST'])
def jpg_to_pdf():
    """Convert JPG/JPEG images to PDF with options - iLovePDF style"""
    tmpdir = None
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No files selected'}), 400
        
        # Get options from form data
        orientation = request.form.get('orientation', 'portrait')  # 'portrait' or 'landscape'
        page_size = request.form.get('pageSize', 'fit')  # 'fit', 'A4', 'Letter', etc.
        margin = request.form.get('margin', 'no')  # 'no', 'small', 'big'
        merge_all = request.form.get('mergeAll', 'false').lower() == 'true'
        
        # Margin settings in points (1 inch = 72 points)
        margin_map = {
            'no': 0,
            'small': 36,  # 0.5 inch
            'big': 72     # 1 inch
        }
        margin_size = margin_map.get(margin, 0)
        
        # Page size settings (width, height in points)
        page_sizes = {
            'A4': (595, 842),
            'Letter': (612, 792),
            'Legal': (612, 1008)
        }
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        print(f"Converting {len(files)} images to PDF (orientation: {orientation}, merge: {merge_all})")
        
        from PIL import Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, letter, legal, landscape, portrait
        from reportlab.lib.utils import ImageReader
        
        if merge_all:
            # Merge all images into one PDF
            pdf_filename = "merged_images.pdf"
            pdf_path = os.path.join(tmpdir, pdf_filename)
            c = canvas.Canvas(pdf_path)
            
            for idx, file in enumerate(files):
                filename = secure_filename(file.filename)
                if not allowed_file(filename, {'jpg', 'jpeg', 'png'}):
                    continue
                
                img_path = os.path.join(tmpdir, filename)
                file.save(img_path)
                
                img = Image.open(img_path)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                img_width, img_height = img.size
                
                # Determine page size
                if page_size == 'fit':
                    if orientation == 'landscape':
                        page_w, page_h = max(img_width, img_height), min(img_width, img_height)
                    else:
                        page_w, page_h = min(img_width, img_height), max(img_width, img_height)
                else:
                    base_size = page_sizes.get(page_size, (595, 842))
                    if orientation == 'landscape':
                        page_w, page_h = max(base_size), min(base_size)
                    else:
                        page_w, page_h = min(base_size), max(base_size)
                
                c.setPageSize((page_w, page_h))
                
                # Calculate image position with margin
                available_w = page_w - (2 * margin_size)
                available_h = page_h - (2 * margin_size)
                
                # Scale image to fit
                scale = min(available_w / img_width, available_h / img_height)
                new_w = img_width * scale
                new_h = img_height * scale
                
                # Center image
                x = (page_w - new_w) / 2
                y = (page_h - new_h) / 2
                
                c.drawImage(img_path, x, y, new_w, new_h)
                c.showPage()
                
                print(f"[OK] Added page {idx+1}/{len(files)}: {filename}")
            
            c.save()
            
            # Read and encode
            with open(pdf_path, 'rb') as f:
                pdf_base64 = base64.b64encode(f.read()).decode('utf-8')
            
            print(f"[OK] Created merged PDF with {len(files)} pages")
            print(f"{'='*60}\n")
            
            return jsonify({
                'success': True,
                'count': 1,
                'pdfs': [{
                    'filename': pdf_filename,
                    'data': pdf_base64
                }]
            })
        else:
            # Create separate PDFs
            pdf_data_list = []
            
            for idx, file in enumerate(files):
                filename = secure_filename(file.filename)
                if not allowed_file(filename, {'jpg', 'jpeg', 'png'}):
                    continue
                
                img_path = os.path.join(tmpdir, filename)
                file.save(img_path)
                
                img = Image.open(img_path)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                img_width, img_height = img.size
                
                # Determine page size
                if page_size == 'fit':
                    if orientation == 'landscape':
                        page_w, page_h = max(img_width, img_height), min(img_width, img_height)
                    else:
                        page_w, page_h = min(img_width, img_height), max(img_width, img_height)
                else:
                    base_size = page_sizes.get(page_size, (595, 842))
                    if orientation == 'landscape':
                        page_w, page_h = max(base_size), min(base_size)
                    else:
                        page_w, page_h = min(base_size), max(base_size)
                
                pdf_filename = f"{Path(filename).stem}.pdf"
                pdf_path = os.path.join(tmpdir, pdf_filename)
                c = canvas.Canvas(pdf_path, pagesize=(page_w, page_h))
                
                # Calculate image position with margin
                available_w = page_w - (2 * margin_size)
                available_h = page_h - (2 * margin_size)
                
                # Scale image to fit
                scale = min(available_w / img_width, available_h / img_height)
                new_w = img_width * scale
                new_h = img_height * scale
                
                # Center image
                x = (page_w - new_w) / 2
                y = (page_h - new_h) / 2
                
                c.drawImage(img_path, x, y, new_w, new_h)
                c.save()
                
                # Read and encode
                with open(pdf_path, 'rb') as f:
                    pdf_base64 = base64.b64encode(f.read()).decode('utf-8')
                    pdf_data_list.append({
                        'filename': pdf_filename,
                        'data': pdf_base64
                    })
                
                print(f"[OK] Converted {idx+1}/{len(files)}: {filename} -> {pdf_filename}")
            
            print(f"[OK] Generated {len(pdf_data_list)} PDFs")
            print(f"{'='*60}\n")
            
            return jsonify({
                'success': True,
                'count': len(pdf_data_list),
                'pdfs': pdf_data_list
            })
    
    except Exception as e:
        # Ensure error message is properly encoded
        error_msg = str(e)
        try:
            print(f"Error: {error_msg}")
            import traceback
            traceback.print_exc()
        except UnicodeEncodeError:
            # If printing fails, just log a simple message
            print("Error occurred during JPG to PDF conversion (encoding issue)")
        return jsonify({'error': error_msg}), 500
    finally:
        # Cleanup
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

@app.route('/api/sign-pdf', methods=['POST'])
def sign_pdf():
    """Add signatures to PDF - iLovePDF style"""
    tmpdir = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'pdf'}):
            return jsonify({'error': 'Only PDF files are allowed'}), 400
        
        # Get signatures data from form
        signatures_json = request.form.get('signatures', '[]')
        signatures = json.loads(signatures_json)
        
        if not signatures:
            return jsonify({'error': 'No signatures provided'}), 400
        
        # Create temp directory
        tmpdir = tempfile.mkdtemp()
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        input_pdf_path = os.path.join(tmpdir, filename)
        file.save(input_pdf_path)
        
        print(f"\n{'='*60}")
        print(f"PDF SIGNING - iLovePDF Style")
        print(f"Input: {filename}")
        print(f"Signatures: {len(signatures)}")
        print(f"{'='*60}\n")
        
        from PyPDF2 import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from PIL import Image as PILImage
        import io as iomodule
        
        # Read the original PDF
        reader = PdfReader(input_pdf_path)
        writer = PdfWriter()
        
        # Process each page
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)
            
            # Get signatures for this page
            page_signatures = [sig for sig in signatures if sig.get('page') == page_num]
            
            if page_signatures:
                # Create overlay with signatures
                packet = iomodule.BytesIO()
                can = canvas.Canvas(packet, pagesize=(page_width, page_height))
                
                for sig in page_signatures:
                    try:
                        # Decode base64 signature image
                        sig_data = sig.get('data', '')
                        if sig_data.startswith('data:image'):
                            sig_data = sig_data.split(',')[1]
                        
                        sig_bytes = base64.b64decode(sig_data)
                        sig_image = PILImage.open(iomodule.BytesIO(sig_bytes))
                        
                        # Save signature as temp image
                        sig_path = os.path.join(tmpdir, f'sig_{page_num}_{signatures.index(sig)}.png')
                        sig_image.save(sig_path, 'PNG')
                        
                        # Get position and size
                        x = float(sig.get('x', 100))
                        y = float(sig.get('y', 100))
                        width = float(sig.get('width', 150))
                        height = float(sig.get('height', 50))
                        
                        # Convert coordinates (PDF uses bottom-left origin)
                        y_converted = page_height - y - height
                        
                        # Draw signature on canvas
                        can.drawImage(sig_path, x, y_converted, width, height, mask='auto')
                        
                        print(f"[OK] Added signature to page {page_num + 1} at ({x}, {y})")
                        
                    except Exception as e:
                        print(f"Warning: Could not add signature: {str(e)}")
                        continue
                
                can.save()
                
                # Move to the beginning of the BytesIO buffer
                packet.seek(0)
                
                # Read the overlay PDF
                overlay_pdf = PdfReader(packet)
                overlay_page = overlay_pdf.pages[0]
                
                # Merge overlay with original page
                page.merge_page(overlay_page)
            
            # Add page to writer
            writer.add_page(page)
        
        # Save signed PDF
        output_filename = f"signed_{filename}"
        output_path = os.path.join(tmpdir, output_filename)
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        print(f"[OK] Created signed PDF: {output_filename}")
        print(f"{'='*60}\n")
        
        # Send signed PDF
        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype='application/pdf'
        )
    
    except Exception as e:
        print(f"Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500
    finally:
        # Cleanup
        if tmpdir and os.path.exists(tmpdir):
            try:
                import time
                time.sleep(1)
                shutil.rmtree(tmpdir, ignore_errors=True)
            except:
                pass

@app.route('/api/health', methods=['GET'])
def api_health_check():
    """Check if LibreOffice is available"""
    try:
        result = subprocess.run(
            ['soffice', '--version'],
            capture_output=True,
            text=True,
            timeout=5,
            encoding='utf-8',
            errors='replace'
        )
        
        if result.returncode == 0:
            return jsonify({
                'status': 'healthy',
                'libreoffice': result.stdout.strip()
            })
        else:
            return jsonify({
                'status': 'unhealthy',
                'error': 'LibreOffice not responding'
            }), 500
    except Exception as e:
        return jsonify({
            'status': 'unhealthy',
            'error': str(e)
        }), 500

@app.route('/edit-pdf', methods=['GET'])
def edit_pdf_page():
    """Render the PDF editor page"""
    return render_template('edit_pdf.html')

@app.route('/api/upload-pdf-for-edit', methods=['POST'])
def upload_pdf_for_edit():
    """Upload PDF and return base64 encoded data for editor"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, {'pdf'}):
            return jsonify({'error': 'Only PDF files are allowed'}), 400
        
        # Read PDF and encode to base64
        pdf_data = file.read()
        base64_pdf = base64.b64encode(pdf_data).decode('utf-8')
        
        # Store in session for later download
        session_id = str(uuid.uuid4())
        session[session_id] = {
            'filename': secure_filename(file.filename),
            'original_data': base64_pdf
        }
        
        return jsonify({
            'success': True,
            'pdf_data': base64_pdf,
            'session_id': session_id,
            'filename': file.filename
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/save-edited-pdf', methods=['POST'])
def save_edited_pdf():
    """Save edited PDF from base64 data"""
    try:
        data = request.get_json()
        
        if not data or 'pdf_data' not in data:
            return jsonify({'error': 'No PDF data provided'}), 400
        
        # Decode base64 PDF
        pdf_base64 = data['pdf_data'].split(',')[1] if ',' in data['pdf_data'] else data['pdf_data']
        pdf_bytes = base64.b64decode(pdf_base64)
        
        # Create temp file
        tmpdir = tempfile.mkdtemp()
        filename = data.get('filename', 'edited-document.pdf')
        pdf_path = os.path.join(tmpdir, filename)
        
        with open(pdf_path, 'wb') as f:
            f.write(pdf_bytes)
        
        # Send file
        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf'
        )
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/', methods=['GET', 'HEAD'])
def index():
    """API Root - Health Check"""
    return jsonify({
        'status': 'healthy',
        'service': 'PDFTools Backend API',
        'version': '3.0.0',
        'message': 'Backend is running successfully!',
        'endpoints': {
            'info': '/api/info',
            'health': '/health'
        }
    })

@app.route('/health', methods=['GET', 'HEAD'])
def health_check():
    """Health check endpoint for Render"""
    return jsonify({'status': 'healthy'}), 200

@app.route('/api/info', methods=['GET'])
def api_info():
    return jsonify({
        'service': 'PDF Converter API',
        'version': '3.0.0',
        'quality': 'iLovePDF Professional Grade',
        'endpoints': {
            'pptx-to-pdf': '/api/convert/pptx-to-pdf (POST)',
            'pdf-to-docx': '/api/convert/pdf-to-docx (POST)',
            'docx-to-pdf': '/api/convert/docx-to-pdf (POST)',
            'pdf-to-excel': '/api/convert/pdf-to-excel (POST)',
            'excel-to-pdf': '/api/convert/excel-to-pdf (POST) - Supports .xlsx, .xls, .csv',
            'edit-pdf': '/edit-pdf (GET) - Interactive PDF Editor',
            'health': '/api/health (GET)'
        },
        'features': {
            'excel_to_pdf': 'Microsoft Excel COM automation for pixel-perfect conversion',
            'csv_support': 'Auto-converts CSV to formatted Excel then PDF',
            'bank_statements': 'Optimized for bank statements with logos and wide tables',
            'landscape_auto': 'Automatically uses landscape for tables with 6+ columns',
            'pdf_editor': 'Full-featured PDF editor with text, images, shapes, and annotations'
        }
    })

# ============================================
# SIGN PDF API ENDPOINTS
# ============================================

@app.route('/api/sign/test', methods=['GET'])
def test_sign_api():
    """Test endpoint to verify API is working"""
    return jsonify({
        'success': True,
        'message': 'Sign PDF API is working!',
        'timestamp': str(pd.Timestamp.now())
    })

@app.route('/api/sign/create-text-signature', methods=['POST'])
def create_text_signature():
    """Create signature from text"""
    try:
        print("=== CREATE TEXT SIGNATURE API CALLED ===")
        data = request.json
        print(f"Request data: {data}")
        text = data.get('text', 'Signature')
        print(f"Text to create: {text}")
        
        # Create signature image from text
        img = Image.new('RGBA', (400, 100), (255, 255, 255, 0))
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except Exception as font_error:
            print(f"Arial font not found, using default: {font_error}")
            font = ImageFont.load_default()
        draw.text((20, 20), text, fill=(0, 0, 0, 255), font=font)
        
        # Convert to base64
        buffered = io.BytesIO()
        img.save(buffered, format="PNG")
        sig_b64 = base64.b64encode(buffered.getvalue()).decode()
        
        print("Signature created successfully")
        return jsonify({
            'success': True,
            'signature': sig_b64,
            'message': 'Signature created successfully'
        })
    
    except Exception as e:
        print(f"ERROR in create_text_signature: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/sign/apply-signatures', methods=['POST', 'OPTIONS'])
def apply_signatures_to_pdf():
    """Apply signatures and stamps to PDF - Using sign_pdf_streamlit.py logic"""
    
    # Handle CORS preflight
    if request.method == 'OPTIONS':
        return jsonify({'success': True}), 200
    
    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas as rcanvas
        from reportlab.lib.utils import ImageReader
        
        print("\n" + "="*50)
        print("=== APPLY SIGNATURES API CALLED ===")
        print("="*50)
        
        # Get request data
        data = request.json
        if not data:
            print("ERROR: No JSON data received")
            return jsonify({'success': False, 'error': 'No data received'}), 400
        
        pdf_data = data.get('pdf_data')  # base64
        placements = data.get('placements', [])  # list of placement objects
        
        print(f"PDF data present: {bool(pdf_data)}")
        print(f"PDF data length: {len(pdf_data) if pdf_data else 0}")
        print(f"Number of placements: {len(placements)}")
        print(f"Placements: {placements}")
        
        if not pdf_data:
            print("ERROR: No PDF data provided")
            return jsonify({'success': False, 'error': 'No PDF data provided'}), 400
            
        if not placements:
            print("ERROR: No placements provided")
            return jsonify({'success': False, 'error': 'No placements provided'}), 400
        
        # Decode PDF
        pdf_bytes = base64.b64decode(pdf_data)
        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        pdf_writer = PdfWriter()
        
        # Group placements by page
        placements_by_page = {}
        for p in placements:
            page_idx = p['page']
            if page_idx not in placements_by_page:
                placements_by_page[page_idx] = []
            placements_by_page[page_idx].append(p)
        
        # Process each page
        for i, page in enumerate(pdf_reader.pages):
            page_w = float(page.mediabox.width)
            page_h = float(page.mediabox.height)
            
            if i in placements_by_page:
                # Create overlay
                overlay_stream = io.BytesIO()
                c = rcanvas.Canvas(overlay_stream, pagesize=(page_w, page_h))
                
                for placement in placements_by_page[i]:
                    # Get image data
                    image_data = placement.get('data')
                    is_stamp = placement.get('type') == 'stamp'
                    
                    if not image_data:
                        continue
                    
                    # Remove data:image/png;base64, prefix if present
                    if 'base64,' in image_data:
                        image_data = image_data.split('base64,')[1]
                    
                    print(f"  - Image data length: {len(image_data)}")
                    
                    # Decode image
                    try:
                        img_bytes = base64.b64decode(image_data)
                        print(f"  - Decoded image bytes: {len(img_bytes)}")
                        img = Image.open(io.BytesIO(img_bytes))
                        print(f"  - Image opened successfully: {img.size}")
                    except Exception as img_error:
                        print(f"  - ERROR decoding image: {img_error}")
                        continue
                    
                    # Save to temp file
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                        img.save(tmp.name, 'PNG')
                        tmp_path = tmp.name
                    
                    try:
                        sig_img = Image.open(tmp_path)
                        sig_w, sig_h = sig_img.size
                        
                        # Get target dimensions
                        target_w = placement.get('width', 150)
                        target_h = placement.get('height', 50)
                        
                        # Get position (from React - top-left origin, pixels)
                        react_x = placement.get('x', 0)
                        react_y = placement.get('y', 0)
                        
                        # Convert from React coordinates (top-left, pixels) to PDF coordinates (bottom-left, points)
                        # React Y=0 is top, PDF Y=0 is bottom
                        abs_x = react_x
                        abs_y = page_h - react_y - target_h
                        
                        item_type = "STAMP" if is_stamp else "SIGNATURE"
                        print(f"\nDEBUG: Page {i+1} - {item_type}")
                        print(f"  - React position: x={react_x:.1f}, y={react_y:.1f}")
                        print(f"  - Page size: w={page_w:.1f}, h={page_h:.1f}")
                        print(f"  - Image size: w={target_w:.1f}, h={target_h:.1f}")
                        print(f"  - PDF position: x={abs_x:.1f}, y={abs_y:.1f}")
                        print(f"  - Drawing image from: {tmp_path}")
                        
                        # Draw image
                        try:
                            c.drawImage(
                                ImageReader(tmp_path),
                                abs_x,
                                abs_y,
                                width=target_w,
                                height=target_h,
                                mask='auto'
                            )
                            print(f"  [OK] Image drawn successfully!")
                        except Exception as draw_error:
                            print(f"  [ERROR] ERROR drawing image: {draw_error}")
                            import traceback
                            traceback.print_exc()
                    finally:
                        try:
                            os.unlink(tmp_path)
                        except:
                            pass
                
                c.save()
                overlay_stream.seek(0)
                
                # Merge overlay
                overlay_pdf = PdfReader(overlay_stream)
                page.merge_page(overlay_pdf.pages[0])
            
            pdf_writer.add_page(page)
        
        # Generate output
        output_stream = io.BytesIO()
        pdf_writer.write(output_stream)
        output_stream.seek(0)
        
        # Return as base64
        signed_pdf_b64 = base64.b64encode(output_stream.getvalue()).decode()
        
        print("="*50)
        print("PDF SIGNED SUCCESSFULLY!")
        print(f"Output PDF size: {len(signed_pdf_b64)} characters")
        print("="*50 + "\n")
        
        return jsonify({
            'success': True,
            'signed_pdf': signed_pdf_b64,
            'message': 'Signatures applied successfully',
            'placements_applied': len(placements)
        })
    
    except Exception as e:
        print("\n" + "="*50)
        print(f"ERROR in apply_signatures_to_pdf: {str(e)}")
        print("="*50)
        import traceback
        traceback.print_exc()
        print("="*50 + "\n")
        return jsonify({'success': False, 'error': str(e)}), 500


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))

def add_text_watermark(pdf_bytes, text, font_size, color, opacity, rotation, position, 
                       font_family='Arial', is_bold=True, is_italic=False, is_underline=False, layer='over'):
    """Add text watermark to PDF with formatting options"""
    from pypdf import PdfReader, PdfWriter
    
    reader = PdfReader(io.BytesIO(pdf_bytes))
    writer = PdfWriter()
    
    # Position mapping - adjusted for proper placement
    pos_map = {
        "center": (0, 0),
        "top-left": (-250, 250),
        "top-center": (0, 250),
        "top-right": (250, 250),
        "middle-left": (-250, 0),
        "middle-right": (250, 0),
        "bottom-left": (-250, -250),
        "bottom-center": (0, -250),
        "bottom-right": (250, -250),
        "tile": None
    }
    
    is_tile = position == "tile"
    is_mosaic = position == "mosaic"
    is_diagonal = position == "diagonal"
    
    # Map font families to ReportLab fonts
    font_map = {
        'Arial': 'Helvetica',
        'Times New Roman': 'Times-Roman',
        'Courier New': 'Courier',
        'Georgia': 'Times-Roman',
        'Verdana': 'Helvetica',
        'Comic Sans MS': 'Helvetica',
        'Impact': 'Helvetica-Bold',
        'Trebuchet MS': 'Helvetica',
    }
    
    base_font = font_map.get(font_family, 'Helvetica')
    
    # Apply bold and italic
    if is_bold and is_italic:
        if base_font == 'Times-Roman':
            font_name = 'Times-BoldItalic'
        elif base_font in ['Helvetica', 'Courier']:
            font_name = f"{base_font}-BoldOblique"
        else:
            font_name = f"{base_font}-Bold"
    elif is_bold:
        if base_font == 'Times-Roman':
            font_name = 'Times-Bold'
        else:
            font_name = f"{base_font}-Bold"
    elif is_italic:
        if base_font == 'Times-Roman':
            font_name = 'Times-Italic'
        elif base_font in ['Helvetica', 'Courier']:
            font_name = f"{base_font}-Oblique"
        else:
            font_name = f"{base_font}-Italic"
    else:
        font_name = base_font
    
    for page in reader.pages:
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=(float(page.mediabox.width), float(page.mediabox.height)))
        w, h = float(page.mediabox.width), float(page.mediabox.height)
        
        # Set color and opacity
        r, g, b = hex_to_rgb(color)
        can.setFillColorRGB(r, g, b, alpha=opacity)
        
        try:
            can.setFont(font_name, font_size)
        except:
            can.setFont("Helvetica-Bold", font_size)
        
        # Calculate text width for better positioning
        text_width = can.stringWidth(text, font_name, font_size)
        text_height = font_size
        
        # Calculate position offsets based on page size (like iLovePDF)
        margin_x = text_width / 2 + 20  # Dynamic margin based on text width
        margin_y = text_height / 2 + 20  # Dynamic margin based on text height
        
        pos_map = {
            "center": (w/2, h/2),
            "top-left": (margin_x, h - margin_y),
            "top-center": (w/2, h - margin_y),
            "top-right": (w - margin_x, h - margin_y),
            "middle-left": (margin_x, h/2),
            "middle-right": (w - margin_x, h/2),
            "bottom-left": (margin_x, margin_y),
            "bottom-center": (w/2, margin_y),
            "bottom-right": (w - margin_x, margin_y),
        }
        
        if is_tile:
            spacing_x = 250
            spacing_y = 150
            for x in range(-int(w), int(w*2), spacing_x):
                for y in range(-int(h), int(h*2), spacing_y):
                    can.saveState()
                    can.translate(w/2 + x, h/2 + y)
                    can.rotate(rotation)
                    can.drawCentredString(0, 0, text)
                    if is_underline:
                        text_width = can.stringWidth(text, font_name, font_size)
                        can.line(-text_width/2, -font_size*0.1, text_width/2, -font_size*0.1)
                    can.restoreState()
        elif is_mosaic:
            # Mosaic pattern - 3x3 grid aligned from top to bottom
            margin_top = font_size / 2 + 10  # Small margin at top
            margin_bottom = font_size / 2 + 10  # Small margin at bottom
            usable_height = h - margin_top - margin_bottom
            
            for row in range(3):
                for col in range(3):
                    x_pos = (col + 0.5) * w / 3
                    # Align from top to bottom with minimal margins
                    y_pos = h - margin_top - (row * usable_height / 2)
                    can.saveState()
                    can.translate(x_pos, y_pos)
                    can.rotate(rotation)
                    can.drawCentredString(0, 0, text)
                    if is_underline:
                        text_width = can.stringWidth(text, font_name, font_size)
                        can.line(-text_width/2, -font_size*0.1, text_width/2, -font_size*0.1)
                    can.restoreState()
        elif is_diagonal:
            # Diagonal watermark from bottom-left to top-right
            import math
            diagonal_length = math.sqrt(w**2 + h**2)
            angle = math.degrees(math.atan2(h, w))
            
            can.saveState()
            can.translate(w/2, h/2)
            can.rotate(angle + rotation)
            can.drawCentredString(0, 0, text)
            if is_underline:
                text_width = can.stringWidth(text, font_name, font_size)
                can.line(-text_width/2, -font_size*0.1, text_width/2, -font_size*0.1)
            can.restoreState()
        else:
            # Get position coordinates
            x_pos, y_pos = pos_map.get(position, (w/2, h/2))
            can.saveState()
            can.translate(x_pos, y_pos)
            can.rotate(rotation)
            can.drawCentredString(0, 0, text)
            if is_underline:
                text_width = can.stringWidth(text, font_name, font_size)
                can.line(-text_width/2, -font_size*0.1, text_width/2, -font_size*0.1)
            can.restoreState()
        
        can.showPage()
        can.save()
        packet.seek(0)
        
        watermark = PdfReader(packet).pages[0]
        if layer == 'below':
            # Watermark below content - merge page onto watermark
            new_page = writer.add_page(watermark)
            new_page.merge_page(page)
        else:
            # Watermark over content (default) - merge watermark onto page
            new_page = writer.add_page(page)
            new_page.merge_page(watermark)
    
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()

def add_image_watermark(pdf_bytes, image_bytes, opacity, rotation, position, layer='over'):
    """Add image watermark to PDF"""
    from pypdf import PdfReader, PdfWriter
    
    reader = PdfReader(io.BytesIO(pdf_bytes))
    writer = PdfWriter()
    
    # Load and process image
    img = Image.open(io.BytesIO(image_bytes))
    
    if img.mode != 'RGBA':
        img = img.convert('RGBA')
    
    alpha = img.split()[3]
    alpha = alpha.point(lambda p: int(p * opacity))
    img.putalpha(alpha)
    
    is_tile = position == "tile"
    is_mosaic = position == "mosaic"
    is_diagonal = position == "diagonal"
    
    max_size = 200
    img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
    img_width, img_height = img.size
    
    for page in reader.pages:
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=(float(page.mediabox.width), float(page.mediabox.height)))
        w, h = float(page.mediabox.width), float(page.mediabox.height)
        
        # Calculate position offsets based on image size (like iLovePDF)
        margin_x = img_width / 2 + 20  # Dynamic margin based on image width
        margin_y = img_height / 2 + 20  # Dynamic margin based on image height
        
        pos_map = {
            "center": (w/2, h/2),
            "top-left": (margin_x, h - margin_y),
            "top-center": (w/2, h - margin_y),
            "top-right": (w - margin_x, h - margin_y),
            "middle-left": (margin_x, h/2),
            "middle-right": (w - margin_x, h/2),
            "bottom-left": (margin_x, margin_y),
            "bottom-center": (w/2, margin_y),
            "bottom-right": (w - margin_x, margin_y),
        }
        
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        img_reader = ImageReader(img_buffer)
        
        if is_tile:
            spacing_x = 300
            spacing_y = 250
            for x in range(-int(w), int(w*2), spacing_x):
                for y in range(-int(h), int(h*2), spacing_y):
                    can.saveState()
                    can.translate(w/2 + x, h/2 + y)
                    can.rotate(rotation)
                    can.drawImage(img_reader, -img_width/2, -img_height/2, 
                                width=img_width, height=img_height, mask='auto')
                    can.restoreState()
        elif is_mosaic:
            # Mosaic pattern - 3x3 grid aligned from top to bottom
            margin_top = img_height / 2 + 10  # Small margin at top
            margin_bottom = img_height / 2 + 10  # Small margin at bottom
            usable_height = h - margin_top - margin_bottom
            
            for row in range(3):
                for col in range(3):
                    x_pos = (col + 0.5) * w / 3
                    # Align from top to bottom with minimal margins
                    y_pos = h - margin_top - (row * usable_height / 2)
                    can.saveState()
                    can.translate(x_pos, y_pos)
                    can.rotate(rotation)
                    can.drawImage(img_reader, -img_width/2, -img_height/2,
                                width=img_width, height=img_height, mask='auto')
                    can.restoreState()
        elif is_diagonal:
            # Diagonal watermark from bottom-left to top-right
            import math
            angle = math.degrees(math.atan2(h, w))
            
            can.saveState()
            can.translate(w/2, h/2)
            can.rotate(angle + rotation)
            can.drawImage(img_reader, -img_width/2, -img_height/2,
                        width=img_width, height=img_height, mask='auto')
            can.restoreState()
        else:
            # Get position coordinates
            x_pos, y_pos = pos_map.get(position, (w/2, h/2))
            can.saveState()
            can.translate(x_pos, y_pos)
            can.rotate(rotation)
            can.drawImage(img_reader, -img_width/2, -img_height/2,
                        width=img_width, height=img_height, mask='auto')
            can.restoreState()
        
        can.showPage()
        can.save()
        packet.seek(0)
        
        watermark = PdfReader(packet).pages[0]
        if layer == 'below':
            # Watermark below content - merge page onto watermark
            new_page = writer.add_page(watermark)
            new_page.merge_page(page)
        else:
            # Watermark over content (default) - merge watermark onto page
            new_page = writer.add_page(page)
            new_page.merge_page(watermark)
    
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()

@app.route('/api/watermark/add', methods=['POST'])
def add_watermark():
    """API endpoint to add watermark to PDF - Optimized for speed"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400
        
        pdf_file = request.files['file']
        pdf_bytes = pdf_file.read()
        
        watermark_type = request.form.get('watermarkType', 'text')
        opacity = float(request.form.get('opacity', 0.5))
        rotation = int(request.form.get('rotation', 0))
        position = request.form.get('position', 'center')
        layer = request.form.get('layer', 'over')
        
        print(f"Adding {watermark_type} watermark to PDF: {pdf_file.filename}")
        
        result_bytes = None
        
        if watermark_type == 'text':
            text = request.form.get('text', 'CONFIDENTIAL')
            font_size = int(request.form.get('fontSize', 40))
            color = request.form.get('color', '#000000')
            font_family = request.form.get('fontFamily', 'Arial')
            is_bold = request.form.get('isBold', 'true').lower() == 'true'
            is_italic = request.form.get('isItalic', 'false').lower() == 'true'
            is_underline = request.form.get('isUnderline', 'false').lower() == 'true'
            
            result_bytes = add_text_watermark(
                pdf_bytes, text, font_size, color, opacity, rotation, position,
                font_family, is_bold, is_italic, is_underline, layer
            )
        
        elif watermark_type == 'image':
            if 'watermarkImage' not in request.files:
                return jsonify({"error": "No watermark image provided"}), 400
            
            image_file = request.files['watermarkImage']
            image_bytes = image_file.read()
            
            result_bytes = add_image_watermark(
                pdf_bytes, image_bytes, opacity, rotation, position, layer
            )
        
        else:
            return jsonify({"error": "Invalid watermark type"}), 400
        
        print(f"Watermark added successfully")
        
        original_name = pdf_file.filename.replace('.pdf', '')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f"{original_name}_watermarked_{timestamp}.pdf"
        
        return send_file(
            io.BytesIO(result_bytes),
            mimetype='application/pdf',
            as_attachment=True,
            download_name=output_filename
        )
    
    except Exception as e:
        print(f"Watermark Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/html-to-pdf', methods=['POST'])
def html_to_pdf():
    """Convert HTML to PDF - Optimized for speed"""
    try:
        # Try to import Playwright
        try:
            from playwright.sync_api import sync_playwright
            use_playwright = True
        except ImportError:
            print("WARNING: Playwright not installed. Using basic HTML to PDF")
            from xhtml2pdf import pisa
            use_playwright = False
        
        url = request.form.get('url')
        html_content = request.form.get('html')
        page_size = request.form.get('pageSize', 'A4')
        orientation = request.form.get('orientation', 'portrait')
        # Optimized margins
        margin_top = request.form.get('marginTop', '5')
        margin_bottom = request.form.get('marginBottom', '5')
        margin_left = request.form.get('marginLeft', '5')
        margin_right = request.form.get('marginRight', '5')
        
        if not url and not html_content:
            return jsonify({"error": "No URL or HTML content provided"}), 400
        
        print(f"Converting HTML to PDF: {url if url else 'custom HTML'}")
        
        if use_playwright:
            # Use Playwright for high-quality PDF generation
            
            try:
                # Configure PDF options
                pdf_options = {
                    'format': page_size,
                    'landscape': orientation == 'landscape',
                    'margin': {
                        'top': f'{margin_top}mm',
                        'bottom': f'{margin_bottom}mm',
                        'left': f'{margin_left}mm',
                        'right': f'{margin_right}mm'
                    },
                    'print_background': True,  # Include background colors and images
                    'prefer_css_page_size': False,
                    'scale': 1.0,  # Full scale
                }
                
                # Generate PDF using Playwright
                print("Launching Playwright browser...")
                with sync_playwright() as p:
                    # Launch browser
                    browser = p.chromium.launch(headless=True)
                    print("Browser launched successfully")
                    page = browser.new_page()
                    
                    # Set viewport to ensure proper rendering
                    page.set_viewport_size({"width": 1920, "height": 1080})
                    
                    if url:
                        # Navigate to URL and wait for everything to load
                        print(f"Loading URL: {url}")
                        page.goto(url, wait_until='load', timeout=30000)
                        
                        # Wait for network to be idle (all resources loaded)
                        print("Waiting for network idle...")
                        page.wait_for_load_state('networkidle', timeout=10000)
                        
                        # Scroll to trigger lazy-loaded images
                        print("Scrolling to load lazy images...")
                        try:
                            # Scroll to bottom
                            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
                            page.wait_for_timeout(1000)
                            # Scroll back to top
                            page.evaluate("window.scrollTo(0, 0)")
                            page.wait_for_timeout(500)
                        except Exception as scroll_error:
                            print(f"Scroll error (non-fatal): {scroll_error}")
                        
                        # Wait for all images to load
                        print("Waiting for images to load...")
                        try:
                            page.wait_for_load_state('domcontentloaded')
                            # Give images time to load
                            page.wait_for_timeout(2000)
                        except Exception as img_error:
                            print(f"Image wait error (non-fatal): {img_error}")
                        
                        # Additional wait for dynamic content
                        print("Final wait for dynamic content...")
                        page.wait_for_timeout(1000)
                    else:
                        # Set HTML content
                        print(f"Setting HTML content ({len(html_content)} chars)")
                        page.set_content(html_content, wait_until='networkidle')
                        page.wait_for_timeout(1000)
                    
                    # Emulate print media for better PDF rendering
                    print("Emulating print media...")
                    page.emulate_media(media='print')
                    
                    # Inject CSS to optimize for PDF (hide unnecessary elements)
                    print("Injecting print optimization CSS...")
                    page.add_style_tag(content="""
                        @media print {
                            /* Hide common elements that create extra pages */
                            footer, .footer, #footer,
                            nav, .nav, .navigation,
                            .sidebar, .side-bar,
                            .advertisement, .ad, .ads,
                            .social-share, .share-buttons,
                            .comments, .comment-section,
                            .related-posts, .related-content,
                            .newsletter, .subscribe,
                            .cookie-banner, .cookie-notice,
                            header.site-header,
                            .back-to-top {
                                display: none !important;
                            }
                            
                            /* Optimize page breaks */
                            body {
                                margin: 0 !important;
                                padding: 0 !important;
                            }
                            
                            /* Prevent page breaks in the middle of content */
                            h1, h2, h3, h4, h5, h6 {
                                page-break-after: avoid !important;
                            }
                            
                            img {
                                page-break-inside: avoid !important;
                            }
                            
                        }
                    """)
                    
                    # Execute JavaScript to limit to ~2 pages of content
                    print("Limiting content to 2 pages...")
                    page.evaluate("""
                        () => {
                            // Find all major sections
                            const sections = document.querySelectorAll('section, article, .section, [role="main"] > div');
                            
                            // Keep only first 2-3 sections (usually covers 2 pages)
                            let keepCount = 0;
                            const maxSections = 2; // Keep first 2 major sections
                            
                            sections.forEach((section, index) => {
                                // Check if section has substantial content
                                const hasContent = section.textContent.trim().length > 100;
                                
                                if (hasContent) {
                                    keepCount++;
                                    if (keepCount > maxSections) {
                                        section.style.display = 'none';
                                    }
                                }
                            });
                            
                            // Also hide specific text patterns that indicate extra sections
                            const allElements = document.querySelectorAll('*');
                            allElements.forEach(el => {
                                const text = el.textContent.trim();
                                // Hide sections with these keywords
                                if (text.includes('Download') || 
                                    text.includes('Buy Prince') ||
                                    text.includes('Invoices') ||
                                    text.includes('Samples') && text.length < 50) {
                                    // Only hide if it's a heading/section, not if it's within content
                                    if (el.tagName.match(/^H[1-6]$/) || el.classList.contains('section')) {
                                        const parent = el.closest('section, article, div[class*="section"]');
                                        if (parent) parent.style.display = 'none';
                                    }
                                }
                            });
                        }
                    """)
                    
                    # Generate PDF
                    print(f"Generating PDF with options: {pdf_options}")
                    pdf_bytes = page.pdf(**pdf_options)
                    
                    browser.close()
                    print("Browser closed")
                
                output = io.BytesIO(pdf_bytes)
                pdf_size = len(pdf_bytes)
                
            except Exception as playwright_error:
                print(f"Playwright error: {str(playwright_error)}")
                import traceback
                traceback.print_exc()
                print("Falling back to xhtml2pdf...")
                # Fall back to xhtml2pdf if Playwright fails
                use_playwright = False
        
        if not use_playwright:
            # Fallback to xhtml2pdf (basic rendering)
            print(f"Converting to PDF using xhtml2pdf (basic rendering)...")
            
            # Get HTML
            if url:
                response = requests.get(url, timeout=30, headers={'User-Agent': 'Mozilla/5.0'})
                response.raise_for_status()
                html_string = response.text
            else:
                html_string = html_content
            
            # Simple conversion
            output = io.BytesIO()
            pisa.CreatePDF(html_string, dest=output, encoding='utf-8')
            pdf_size = len(output.getvalue())
        
        output.seek(0)
        
        # Check if PDF was generated
        print(f"Generated PDF size: {pdf_size} bytes")
        
        if pdf_size == 0:
            raise Exception("PDF generation failed - output is empty")
        
        # Generate filename
        if url:
            from urllib.parse import urlparse
            hostname = urlparse(url).hostname or 'webpage'
            filename = f"{hostname}.pdf"
        else:
            filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        print(f"Sending PDF: {filename} ({pdf_size} bytes)")
        
        return send_file(
            output,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=filename
        )
    
    except requests.exceptions.RequestException as e:
        print(f"URL Fetch Error: {str(e)}")
        return jsonify({"error": f"Failed to fetch URL: {str(e)}"}), 400
    except Exception as e:
        print(f"HTML to PDF Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/pdf/rotate', methods=['POST'])
def rotate_pdf():
    """API endpoint to rotate PDF pages - Optimized for speed"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400
        
        pdf_file = request.files['file']
        if not pdf_file.filename:
            return jsonify({"error": "No file selected"}), 400
            
        pdf_bytes = pdf_file.read()
        
        if len(pdf_bytes) == 0:
            return jsonify({"error": "Empty file provided"}), 400
        
        # Get rotation parameters
        rotation_angle = int(request.form.get('rotation', 90))
        direction = request.form.get('direction', 'right')
        
        print(f"Rotating PDF: {pdf_file.filename} ({rotation_angle}° {direction})")
        
        # Calculate actual rotation
        if direction == 'left':
            rotation_angle = 360 - rotation_angle
        
        from pypdf import PdfReader, PdfWriter
        
        reader = PdfReader(io.BytesIO(pdf_bytes))
        writer = PdfWriter()
        
        # Rotate all pages efficiently
        for page in reader.pages:
            page.rotate(rotation_angle)
            writer.add_page(page)
        
        # Generate output
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        print(f"PDF rotated successfully: {len(reader.pages)} pages")
        
        original_name = pdf_file.filename.replace('.pdf', '')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f"{original_name}_rotated_{timestamp}.pdf"
        
        return send_file(
            output,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=output_filename
        )
    
    except Exception as e:
        print(f"Rotation Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/pdf/unlock', methods=['POST'])
def unlock_pdf():
    """API endpoint to unlock/decrypt password-protected PDF files - Optimized"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400
        
        pdf_file = request.files['file']
        if not pdf_file.filename:
            return jsonify({"error": "No file selected"}), 400
            
        pdf_bytes = pdf_file.read()
        
        if len(pdf_bytes) == 0:
            return jsonify({"error": "Empty file provided"}), 400
        
        # Get password from request (if provided)
        password = request.form.get('password', '').strip()
        
        print(f"Unlocking PDF: {pdf_file.filename}")
        
        from pypdf import PdfReader, PdfWriter
        
        # Try to read the PDF with password
        reader = PdfReader(io.BytesIO(pdf_bytes))
        
        # Check if PDF is encrypted
        if reader.is_encrypted:
            print(f"PDF is encrypted, attempting to decrypt")
            
            # Try to decrypt with provided password
            if password:
                try:
                    decrypt_result = reader.decrypt(password)
                    print(f"Decrypt result: {decrypt_result}")
                    
                    # decrypt_result: 0 = failed, 1 = user password, 2 = owner password
                    if decrypt_result == 0:
                        return jsonify({
                            "error": "Incorrect password. Please try again.",
                            "needsPassword": True
                        }), 400
                    
                    print(f"Successfully decrypted with password (result: {decrypt_result})")
                except Exception as decrypt_err:
                    print(f"Decryption exception: {str(decrypt_err)}")
                    return jsonify({
                        "error": "Incorrect password. Please try again.",
                        "needsPassword": True
                    }), 400
            else:
                return jsonify({
                    "error": "This PDF is password-protected. Please provide the password.",
                    "needsPassword": True
                }), 400
        else:
            print("PDF is not encrypted, proceeding with unlock")
        
        # Create a new PDF without encryption
        writer = PdfWriter()
        
        # Copy all pages to new writer (without encryption)
        for page in reader.pages:
            writer.add_page(page)
        
        # Copy metadata if available
        if reader.metadata:
            writer.add_metadata(reader.metadata)
        
        # Generate output
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        original_name = pdf_file.filename.replace('.pdf', '')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f"{original_name}_unlocked_{timestamp}.pdf"
        
        print(f"Successfully unlocked PDF: {output_filename}")
        
        return send_file(
            output,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=output_filename
        )
    
    except Exception as e:
        print(f"Unlock PDF Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"Failed to process PDF: {str(e)}"}), 500

@app.route('/api/pdf/protect', methods=['POST'])
def protect_pdf():
    """API endpoint to protect/encrypt PDF files with password - Optimized"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400
        
        pdf_file = request.files['file']
        if not pdf_file.filename:
            return jsonify({"error": "No file selected"}), 400
            
        pdf_bytes = pdf_file.read()
        
        if len(pdf_bytes) == 0:
            return jsonify({"error": "Empty file provided"}), 400
        
        # Get password from request
        password = request.form.get('password', '').strip()
        
        if not password:
            return jsonify({"error": "Password is required"}), 400
        
        if len(password) < 6:
            return jsonify({"error": "Password must be at least 6 characters long"}), 400
        
        print(f"Protecting PDF with password: {pdf_file.filename}")
        
        from pypdf import PdfReader, PdfWriter
        
        # Read the PDF
        reader = PdfReader(io.BytesIO(pdf_bytes))
        writer = PdfWriter()
        
        # Copy all pages to new writer efficiently
        for page in reader.pages:
            writer.add_page(page)
        
        # Copy metadata if available
        if reader.metadata:
            writer.add_metadata(reader.metadata)
        
        # Encrypt the PDF with AES-256 encryption
        writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")
        
        print(f"PDF encrypted successfully with {len(reader.pages)} pages")
        
        # Generate output
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        original_name = pdf_file.filename.replace('.pdf', '')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f"{original_name}_protected_{timestamp}.pdf"
        
        print(f"Successfully protected PDF: {output_filename}")
        
        return send_file(
            output,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=output_filename
        )
    
    except Exception as e:
        print(f"Protect PDF Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"Failed to protect PDF: {str(e)}"}), 500

if __name__ == '__main__':
    # Get port from environment variable (for Render deployment) or use 5000 for local
    port = int(os.environ.get('PORT', 5000))
    
    print("Starting PDFTools Backend API")
    print(f"Server: http://0.0.0.0:{port}")
    print(f"Health Check: http://0.0.0.0:{port}/health")
    print(f"API Info: http://0.0.0.0:{port}/api/info")
    print("\nAvailable Endpoints:")
    print("  - PowerPoint to PDF: /api/convert/pptx-to-pdf")
    print("  - Sign PDF: /api/sign/apply-signatures")
    print("  - Watermark PDF: /api/watermark/add")
    print("  - Rotate PDF: /api/pdf/rotate")
    print("  - Unlock PDF: /api/pdf/unlock")
    print("  - Protect PDF: /api/pdf/protect")
    
    # Disable debug mode in production
    debug_mode = os.environ.get('FLASK_ENV') != 'production'
    app.run(host='0.0.0.0', port=port, debug=debug_mode)

