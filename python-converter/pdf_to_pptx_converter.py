#!/usr/bin/env python3
"""
PDF to PPTX Converter
Converts PDF pages to images and creates PowerPoint slides
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from PIL import Image
import tempfile

def convert_pdf_to_pptx_images(pdf_path, output_dir):
    """
    Convert PDF to PPTX by converting each page to an image
    
    Args:
        pdf_path: Path to input PDF file
        output_dir: Directory for output PPTX file
    
    Returns:
        Path to converted PPTX file
    """
    try:
        # Convert PDF pages to images
        print(f"Converting PDF pages to images...")
        images = convert_from_path(
            pdf_path,
            dpi=300,  # High quality
            fmt='png'
        )
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions to match standard 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        print(f"Creating PowerPoint with {len(images)} slides...")
        
        # Add each image as a slide
        for i, image in enumerate(images):
            print(f"  Adding slide {i+1}/{len(images)}...")
            
            # Save image temporarily
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                image.save(tmp_img.name, 'PNG')
                tmp_img_path = tmp_img.name
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate dimensions to fit slide
            img_width, img_height = image.size
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Calculate scaling to fit
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale = min(width_ratio, height_ratio)
            
            # Calculate centered position
            new_width = int(img_width * scale)
            new_height = int(img_height * scale)
            left = (slide_width - new_width) // 2
            top = (slide_height - new_height) // 2
            
            # Add image to slide
            slide.shapes.add_picture(
                tmp_img_path,
                left, top,
                width=new_width,
                height=new_height
            )
            
            # Clean up temp image
            try:
                os.unlink(tmp_img_path)
            except:
                pass
        
        # Save PPTX
        pptx_name = Path(pdf_path).stem + '.pptx'
        pptx_path = Path(output_dir) / pptx_name
        prs.save(str(pptx_path))
        
        print(f"✅ Success! Created: {pptx_path}")
        return pptx_path
        
    except Exception as e:
        raise RuntimeError(f"PDF to PPTX conversion failed: {str(e)}")

if __name__ == '__main__':
    import sys
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx_converter.py <input.pdf> <output_dir>")
        sys.exit(1)
    
    pdf_path = sys.argv[1]
    output_dir = sys.argv[2]
    
    try:
        pptx_path = convert_pdf_to_pptx_images(pdf_path, output_dir)
        print(f"Success! Created: {pptx_path}")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
