"""
Fast PDF to PowerPoint converter optimized for cloud deployments (Render, Vercel)
Uses PyMuPDF (fitz) for faster image extraction
Optimized for speed over quality
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF


def pdf_to_pptx_fast(pdf_path, output_dir=None, dpi=150):
    """
    Fast PDF to PPTX conversion using PyMuPDF
    Optimized for cloud deployments - 3x faster than pdf2image
    
    Args:
        pdf_path: Path to input PDF
        output_dir: Output directory (default: same as input)
        dpi: Image resolution (default: 150 for speed, use 200 for quality)
    """
    if output_dir is None:
        output_dir = os.path.dirname(pdf_path) or '.'
    
    print(f"\n{'='*60}")
    print(f"FAST PDF TO PPTX CONVERSION (Cloud-Optimized)")
    print(f"Input: {pdf_path}")
    print(f"DPI: {dpi} (optimized for speed)")
    print(f"{'='*60}\n")
    
    # Open PDF with PyMuPDF (much faster than pdf2image)
    doc = fitz.open(pdf_path)
    page_count = len(doc)
    
    print(f"PDF has {page_count} pages")
    
    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    
    # Convert each page to image and add to PPTX
    for page_num in range(page_count):
        page = doc[page_num]
        
        # Calculate zoom for desired DPI
        # PyMuPDF default is 72 DPI
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        
        # Render page to pixmap (image)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        
        # Save as PNG in memory
        img_data = pix.tobytes("png")
        
        # Create temporary file for image
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
            tmp_img.write(img_data)
            tmp_img_path = tmp_img.name
        
        # Add slide with image as full background
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            tmp_img_path,
            0,  # left = 0
            0,  # top = 0
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # Clean up temp file
        try:
            os.unlink(tmp_img_path)
        except:
            pass
        
        print(f"Added slide {page_num + 1}/{page_count}")
    
    doc.close()
    
    # Save PPTX
    pptx_name = Path(pdf_path).stem + '.pptx'
    pptx_path = Path(output_dir) / pptx_name
    prs.save(str(pptx_path))
    
    print(f"\n[OK] PowerPoint created: {pptx_path}")
    print(f"Size: {os.path.getsize(pptx_path) / 1024:.1f} KB")
    print(f"{'='*60}\n")
    
    return str(pptx_path)


def pdf_to_pptx_ultra_fast(pdf_path, output_dir=None):
    """
    Ultra-fast conversion with minimal quality (DPI=100)
    Use for large PDFs or when speed is critical
    """
    return pdf_to_pptx_fast(pdf_path, output_dir, dpi=100)


def pdf_to_pptx_balanced(pdf_path, output_dir=None):
    """
    Balanced conversion (DPI=150) - default
    Good quality with fast speed
    """
    return pdf_to_pptx_fast(pdf_path, output_dir, dpi=150)


def pdf_to_pptx_quality(pdf_path, output_dir=None):
    """
    Higher quality conversion (DPI=200)
    Slower but better image quality
    """
    return pdf_to_pptx_fast(pdf_path, output_dir, dpi=200)


def convert_pdf_to_pptx_optimized(pdf_path, output_dir=None, quality='balanced'):
    """
    Main conversion function with quality options
    
    Args:
        pdf_path: Path to input PDF
        output_dir: Output directory
        quality: 'ultra_fast' (100 DPI), 'balanced' (150 DPI), 'quality' (200 DPI)
    """
    quality_map = {
        'ultra_fast': 100,
        'balanced': 150,
        'quality': 200
    }
    
    dpi = quality_map.get(quality, 150)
    
    try:
        return pdf_to_pptx_fast(pdf_path, output_dir, dpi=dpi)
    except Exception as e:
        print(f"Fast converter failed: {e}")
        raise


if __name__ == "__main__":
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
        output_dir = sys.argv[2] if len(sys.argv) > 2 else None
        quality = sys.argv[3] if len(sys.argv) > 3 else 'balanced'
        
        result = convert_pdf_to_pptx_optimized(input_file, output_dir, quality)
        print(f"Result: {result}")
    else:
        print("Usage: python pdf_to_pptx_fast.py <input.pdf> [output_dir] [quality]")
        print("Quality options: ultra_fast, balanced (default), quality")
